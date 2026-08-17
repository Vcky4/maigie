"""
Document Generation Service.

Generates downloadable documents (PDF, DOCX) from HTML content.
Uses WeasyPrint for PDF rendering and htmldocx for DOCX conversion.
"""

from __future__ import annotations

import io
import logging
import re
import uuid
from datetime import UTC, datetime
from typing import Any

logger = logging.getLogger(__name__)

# Content type mappings
CONTENT_TYPES = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

# Max content length to prevent abuse (100k chars ~= 50 pages)
MAX_CONTENT_LENGTH = 100_000

# CSS styles for different document styles
_ACADEMIC_CSS = """
@page {
    size: A4;
    margin: 2.5cm 2cm;
    @bottom-center {
        content: counter(page) " / " counter(pages);
        font-family: 'Times New Roman', Times, serif;
        font-size: 9pt;
        color: #888;
    }
}
@page :first {
    @bottom-center { content: none; }
}
body {
    font-family: 'Times New Roman', Times, Georgia, serif;
    font-size: 12pt;
    line-height: 1.6;
    color: #1a1a1a;
    orphans: 3;
    widows: 3;
    hyphens: auto;
}
h1 {
    font-size: 22pt;
    margin-top: 1.5em;
    margin-bottom: 0.5em;
    color: #111;
    page-break-after: avoid;
}
h2 {
    font-size: 16pt;
    margin-top: 1.3em;
    margin-bottom: 0.4em;
    color: #222;
    page-break-after: avoid;
}
h3 {
    font-size: 13pt;
    margin-top: 1.1em;
    margin-bottom: 0.3em;
    color: #333;
    page-break-after: avoid;
}
h4 { font-size: 12pt; margin-top: 1em; margin-bottom: 0.3em; font-style: italic; page-break-after: avoid; }
p { margin-bottom: 0.8em; text-align: justify; }
ul, ol { margin-bottom: 0.8em; padding-left: 2em; }
li { margin-bottom: 0.3em; }
blockquote {
    border-left: 3px solid #ccc;
    margin-left: 0;
    padding-left: 1em;
    color: #555;
    font-style: italic;
}
code {
    font-family: 'Courier New', monospace;
    font-size: 10pt;
    background: #f5f5f5;
    padding: 1px 4px;
}
pre {
    background: #f5f5f5;
    padding: 12px;
    font-size: 10pt;
    line-height: 1.4;
    white-space: pre-wrap;
    word-wrap: break-word;
}
pre code { background: none; padding: 0; }
table {
    width: 100%;
    border-collapse: collapse;
    margin: 1em 0;
    font-size: 11pt;
}
th, td {
    border: 1px solid #ccc;
    padding: 8px 10px;
    text-align: left;
}
th { background: #f0f0f0; font-weight: bold; }
hr { border: none; border-top: 1px solid #ddd; margin: 2em 0; }
"""

_REPORT_CSS = """
@page {
    size: A4;
    margin: 2cm;
    @top-right {
        content: string(doc-title);
        font-family: Helvetica, Arial, sans-serif;
        font-size: 9pt;
        color: #888;
    }
    @bottom-center {
        content: "Page " counter(page) " of " counter(pages);
        font-family: Helvetica, Arial, sans-serif;
        font-size: 9pt;
        color: #888;
    }
}
@page :first {
    @top-right { content: none; }
    @bottom-center { content: none; }
}
body {
    font-family: Helvetica, Arial, sans-serif;
    font-size: 11pt;
    line-height: 1.5;
    color: #2d2d2d;
    orphans: 3;
    widows: 3;
}
h1 {
    font-size: 24pt;
    margin-top: 1.2em;
    margin-bottom: 0.4em;
    color: #1a1a1a;
    font-weight: bold;
    string-set: doc-title content();
    page-break-after: avoid;
}
h2 {
    font-size: 16pt;
    margin-top: 1.2em;
    margin-bottom: 0.4em;
    color: #333;
    border-bottom: 1px solid #eee;
    padding-bottom: 4px;
    page-break-after: avoid;
}
h3 { font-size: 13pt; margin-top: 1em; margin-bottom: 0.3em; color: #444; page-break-after: avoid; }
h4 { font-size: 11pt; margin-top: 0.8em; margin-bottom: 0.3em; color: #555; font-weight: bold; page-break-after: avoid; }
p { margin-bottom: 0.7em; }
ul, ol { margin-bottom: 0.7em; padding-left: 1.8em; }
li { margin-bottom: 0.2em; }
blockquote {
    border-left: 4px solid #4f46e5;
    margin-left: 0;
    padding-left: 1em;
    color: #555;
}
code {
    font-family: 'Courier New', monospace;
    font-size: 9.5pt;
    background: #f8f8f8;
    padding: 2px 5px;
}
pre {
    background: #f8f8f8;
    padding: 14px;
    font-size: 9.5pt;
    line-height: 1.4;
    border: 1px solid #e8e8e8;
    white-space: pre-wrap;
    word-wrap: break-word;
}
pre code { background: none; padding: 0; border: none; }
table {
    width: 100%;
    border-collapse: collapse;
    margin: 1em 0;
    font-size: 10pt;
}
th, td {
    border: 1px solid #ddd;
    padding: 8px 12px;
    text-align: left;
}
th { background: #f5f5f5; font-weight: bold; }
hr { border: none; border-top: 2px solid #eee; margin: 1.5em 0; }
"""

_MINIMAL_CSS = """
@page {
    size: A4;
    margin: 1.5cm;
}
body {
    font-family: Helvetica, Arial, sans-serif;
    font-size: 10.5pt;
    line-height: 1.45;
    color: #333;
}
h1 { font-size: 18pt; margin-top: 1em; margin-bottom: 0.3em; }
h2 { font-size: 14pt; margin-top: 0.9em; margin-bottom: 0.3em; }
h3 { font-size: 12pt; margin-top: 0.8em; margin-bottom: 0.2em; }
p { margin-bottom: 0.5em; }
ul, ol { margin-bottom: 0.5em; padding-left: 1.5em; }
li { margin-bottom: 0.15em; }
code { font-family: monospace; font-size: 9.5pt; background: #f5f5f5; padding: 1px 3px; }
pre { background: #f5f5f5; padding: 10px; font-size: 9pt; white-space: pre-wrap; word-wrap: break-word; }
pre code { background: none; padding: 0; }
table { width: 100%; border-collapse: collapse; margin: 0.8em 0; font-size: 9.5pt; }
th, td { border: 1px solid #ddd; padding: 6px 8px; }
th { background: #f0f0f0; font-weight: bold; }
hr { border: none; border-top: 1px solid #ddd; margin: 1em 0; }
"""

_STYLE_CSS = {
    "academic": _ACADEMIC_CSS,
    "report": _REPORT_CSS,
    "minimal": _MINIMAL_CSS,
}


class DocumentGenerationService:
    """Generates PDF and DOCX documents from HTML content."""

    def __init__(self):
        pass

    async def generate_document(
        self,
        format: str,
        title: str,
        content: str,
        style: str = "academic",
        user_id: str | None = None,
        doc_type: str | None = None,
    ) -> dict[str, Any]:
        """
        Generate a document in the specified format.

        Args:
            format: Document format ("pdf" or "docx")
            title: Document title
            content: HTML content to render (also handles markdown via conversion)
            style: Document style ("academic", "report", "minimal")
            user_id: User ID for path namespacing

        Returns:
            dict with keys: filename, url, size, format, content_type
        """
        if len(content) > MAX_CONTENT_LENGTH:
            content = content[:MAX_CONTENT_LENGTH]
            logger.warning(f"Content truncated to {MAX_CONTENT_LENGTH} chars for user {user_id}")

        format = _normalize_format(format)
        is_presentation = (doc_type or "").strip().lower() == "presentation"

        # Presentation as PDF gets a first-class HTML slide renderer.
        if is_presentation and format == "pdf":
            slides = self._extract_slides_from_content(content, title)
            doc_bytes = render_presentation_pdf(title, slides, theme="indigo")
            preview_html = build_presentation_html(title, slides, theme="indigo")
        elif format == "pptx":
            doc_bytes = self._generate_pptx(title, content, style)
            preview_html = self._build_pptx_preview_html(title, content, style)
        else:
            # Normalize content: if it's markdown, convert to HTML
            content = self._ensure_html(content)
            if format == "pdf":
                doc_bytes = self._generate_pdf(title, content, style)
            else:
                doc_bytes = self._generate_docx(title, content, style)
            preview_html = self._build_full_html(title, content, style)

        # Generate a unique filename
        safe_title = re.sub(r"[^\w\s-]", "", title)[:50].strip().replace(" ", "_")
        timestamp = datetime.now(UTC).strftime("%Y%m%d_%H%M%S")
        short_id = uuid.uuid4().hex[:6]
        filename = f"{safe_title}_{timestamp}_{short_id}.{format}"

        # Upload to storage
        from src.shared.infrastructure.storage_service import storage_service

        storage_path = f"generated-docs/{user_id or 'anonymous'}"
        upload_result = await self._upload_bytes(storage_service, doc_bytes, filename, storage_path)

        # Also upload the styled HTML for in-app preview
        html_filename = f"{safe_title}_{timestamp}_{short_id}.html"
        preview_result = await self._upload_bytes(
            storage_service, preview_html.encode("utf-8"), html_filename, storage_path
        )

        return {
            "filename": filename,
            "url": upload_result["url"],
            "size": upload_result["size"],
            "format": format,
            "content_type": CONTENT_TYPES[format],
            "title": title,
            "preview_url": preview_result["url"],
        }

    def _ensure_html(self, content: str) -> str:
        """If content is markdown, convert to HTML. If already HTML, return as-is."""
        # Detect if content is already HTML (has block-level HTML tags)
        if re.search(r"<(h[1-6]|p|ul|ol|li|table|div|pre|blockquote)\b", content, re.IGNORECASE):
            return content

        # Content is markdown — convert to HTML
        return self._markdown_to_html(content)

    def _markdown_to_html(self, md: str) -> str:
        """Convert markdown to HTML for document rendering."""
        lines = md.split("\n")
        html_parts: list[str] = []
        in_code_block = False
        code_buffer: list[str] = []
        in_list = False
        list_type = ""
        paragraph_buffer: list[str] = []

        def flush_paragraph():
            if paragraph_buffer:
                text = " ".join(paragraph_buffer)
                text = self._inline_md_to_html(text)
                html_parts.append(f"<p>{text}</p>")
                paragraph_buffer.clear()

        def flush_list():
            nonlocal in_list, list_type
            if in_list:
                html_parts.append(f"</{list_type}>")
                in_list = False
                list_type = ""

        for line in lines:
            # Code blocks
            if line.strip().startswith("```"):
                if in_code_block:
                    html_parts.append(
                        f"<pre><code>{self._escape_html(chr(10).join(code_buffer))}</code></pre>"
                    )
                    code_buffer = []
                    in_code_block = False
                else:
                    flush_paragraph()
                    flush_list()
                    in_code_block = True
                continue

            if in_code_block:
                code_buffer.append(line)
                continue

            # Headings
            heading_match = re.match(r"^(#{1,6})\s+(.+)$", line)
            if heading_match:
                flush_paragraph()
                flush_list()
                level = len(heading_match.group(1))
                text = self._inline_md_to_html(heading_match.group(2))
                html_parts.append(f"<h{level}>{text}</h{level}>")
                continue

            # Horizontal rule
            if re.match(r"^(\-{3,}|\*{3,}|_{3,})\s*$", line):
                flush_paragraph()
                flush_list()
                html_parts.append("<hr>")
                continue

            # Bullet lists
            bullet_match = re.match(r"^(\s*)([-*+])\s+(.+)$", line)
            if bullet_match:
                flush_paragraph()
                if not in_list or list_type != "ul":
                    flush_list()
                    html_parts.append("<ul>")
                    in_list = True
                    list_type = "ul"
                text = self._inline_md_to_html(bullet_match.group(3))
                html_parts.append(f"<li>{text}</li>")
                continue

            # Numbered lists
            numbered_match = re.match(r"^(\s*)\d+\.\s+(.+)$", line)
            if numbered_match:
                flush_paragraph()
                if not in_list or list_type != "ol":
                    flush_list()
                    html_parts.append("<ol>")
                    in_list = True
                    list_type = "ol"
                text = self._inline_md_to_html(numbered_match.group(2))
                html_parts.append(f"<li>{text}</li>")
                continue

            # Blockquotes
            if line.startswith("> "):
                flush_paragraph()
                flush_list()
                text = self._inline_md_to_html(line[2:])
                html_parts.append(f"<blockquote><p>{text}</p></blockquote>")
                continue

            # Empty line
            if not line.strip():
                flush_paragraph()
                flush_list()
                continue

            # Regular text — accumulate for paragraph
            flush_list()
            paragraph_buffer.append(line.strip())

        # Flush remaining
        flush_paragraph()
        flush_list()
        if in_code_block and code_buffer:
            html_parts.append(
                f"<pre><code>{self._escape_html(chr(10).join(code_buffer))}</code></pre>"
            )

        return "\n".join(html_parts)

    def _inline_md_to_html(self, text: str) -> str:
        """Convert inline markdown (bold, italic, code, links) to HTML."""
        # Bold
        text = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", text)
        text = re.sub(r"__(.+?)__", r"<b>\1</b>", text)
        # Italic
        text = re.sub(r"\*(.+?)\*", r"<i>\1</i>", text)
        text = re.sub(r"_(.+?)_", r"<i>\1</i>", text)
        # Inline code
        text = re.sub(r"`(.+?)`", r"<code>\1</code>", text)
        # Links
        text = re.sub(r"\[(.+?)\]\((.+?)\)", r'<a href="\2">\1</a>', text)
        return text

    def _escape_html(self, text: str) -> str:
        """Escape HTML special characters."""
        return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    def _build_full_html(self, title: str, body_html: str, style: str) -> str:
        """Wrap body HTML with full document structure and CSS."""
        css = _STYLE_CSS.get(style, _ACADEMIC_CSS)
        return f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<title>{self._escape_html(title)}</title>
<style>{css}</style>
</head>
<body>
{body_html}
</body>
</html>"""

    def _generate_pdf(self, title: str, content: str, style: str) -> bytes:
        """
        Generate a PDF document from HTML content.

        Prefers WeasyPrint for print-quality output (proper fonts, page control,
        widow/orphan handling, real CSS support). Falls back to xhtml2pdf when
        WeasyPrint's native libraries (Pango, Cairo) are unavailable.
        """
        full_html = self._build_full_html(title, content, style)

        try:
            from weasyprint import HTML as _WeasyHTML

            return _WeasyHTML(string=full_html).write_pdf()
        except (ImportError, OSError) as e:
            # OSError covers "cannot load library libgobject-2.0-0" on macOS
            # when Homebrew Pango is not on DYLD_FALLBACK_LIBRARY_PATH.
            logger.warning(
                f"WeasyPrint unavailable ({e.__class__.__name__}: {e}). "
                "Falling back to xhtml2pdf. See docs to enable WeasyPrint."
            )

        from xhtml2pdf import pisa

        buffer = io.BytesIO()
        pisa_status = pisa.CreatePDF(io.StringIO(full_html), dest=buffer)
        if pisa_status.err:
            logger.error(f"xhtml2pdf error count: {pisa_status.err}")

        buffer.seek(0)
        return buffer.getvalue()

    def _generate_docx(self, title: str, content: str, style: str) -> bytes:
        """Generate a DOCX document from HTML content using htmldocx."""
        from docx import Document
        from docx.shared import Pt
        from htmldocx import HtmlToDocx

        doc = Document()
        parser = HtmlToDocx()

        # Add the HTML content
        parser.add_html_to_document(content, doc)

        # Save to bytes
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    def _generate_pptx(self, title: str, content: str, style: str) -> bytes:
        """Generate a richly styled PPTX from content with tables, shapes, and layouts."""
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Color palette
        PRIMARY = RGBColor(0x4F, 0x46, 0xE5)
        DARK_BG = RGBColor(0x1E, 0x1B, 0x4B)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
        TEXT_DARK = RGBColor(0x1E, 0x29, 0x3B)
        TEXT_GRAY = RGBColor(0x64, 0x74, 0x8B)
        ACCENT2 = RGBColor(0x7C, 0x3A, 0xED)
        ACCENT3 = RGBColor(0x06, 0xB6, 0xD4)

        slides_data = self._extract_slides_from_content(content, title)

        # Filter out slides with only a title and no real content
        # Keep: first slide (title slide), slides with bullets, slides with tables
        filtered = []
        for idx, s in enumerate(slides_data):
            if idx == 0:
                filtered.append(s)
            elif s.get("bullets") or s.get("table"):
                filtered.append(s)
        slides_data = filtered if filtered else slides_data[:1]

        for i, slide_data in enumerate(slides_data):
            slide_title = slide_data.get("title", f"Slide {i + 1}")
            bullets = slide_data.get("bullets", [])
            subtitle = slide_data.get("subtitle", "")
            table_data = slide_data.get("table")

            layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(layout)

            if i == 0:
                # ═══ TITLE SLIDE ═══
                bg = slide.background.fill
                bg.solid()
                bg.fore_color.rgb = DARK_BG

                # Top accent bar
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.12)
                )
                bar.fill.solid()
                bar.fill.fore_color.rgb = PRIMARY
                bar.line.fill.background()

                # Left accent stripe
                acc = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(0.08), Inches(2.5)
                )
                acc.fill.solid()
                acc.fill.fore_color.rgb = PRIMARY
                acc.line.fill.background()

                # Title
                tb = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10), Inches(1.8))
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_title
                p.font.size = Pt(48)
                p.font.bold = True
                p.font.color.rgb = WHITE

                # Subtitle
                sub_text = subtitle or (bullets[0] if bullets else "")
                if sub_text:
                    tb2 = slide.shapes.add_textbox(Inches(1.2), Inches(4.2), Inches(10), Inches(1))
                    tf2 = tb2.text_frame
                    tf2.word_wrap = True
                    p2 = tf2.paragraphs[0]
                    p2.text = sub_text
                    p2.font.size = Pt(20)
                    p2.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)

                # Decorative circles
                for cx, clr in [(Inches(10), ACCENT2), (Inches(11.2), ACCENT3)]:
                    c = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL, cx, Inches(5.5), Inches(1.5), Inches(1.5)
                    )
                    c.fill.solid()
                    c.fill.fore_color.rgb = clr
                    c.line.fill.background()

            elif table_data and isinstance(table_data, dict):
                # ═══ TABLE SLIDE ═══
                self._pptx_add_header(slide, slide_title, prs.slide_width, PRIMARY, WHITE)
                headers = table_data.get("headers", [])
                rows = table_data.get("rows", [])
                if headers and rows:
                    n_rows = len(rows) + 1
                    n_cols = len(headers)
                    tbl_w = min(prs.slide_width - Inches(1.6), Inches(n_cols * 2.5))
                    tbl_left = (prs.slide_width - tbl_w) // 2
                    shape = slide.shapes.add_table(
                        n_rows, n_cols, tbl_left, Inches(2.0), tbl_w, Inches(4.5)
                    )
                    tbl = shape.table
                    tbl.first_row = True
                    for ci, h in enumerate(headers):
                        cell = tbl.cell(0, ci)
                        cell.text = str(h)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = PRIMARY
                        for para in cell.text_frame.paragraphs:
                            para.font.color.rgb = WHITE
                            para.font.size = Pt(14)
                            para.font.bold = True
                    for ri, row in enumerate(rows):
                        for ci, val in enumerate(row[:n_cols]):
                            cell = tbl.cell(ri + 1, ci)
                            cell.text = str(val)
                            for para in cell.text_frame.paragraphs:
                                para.font.size = Pt(13)
                                para.font.color.rgb = TEXT_DARK
                            if ri % 2 == 1:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                # ═══ CONTENT SLIDE ═══
                self._pptx_add_header(slide, slide_title, prs.slide_width, PRIMARY, WHITE)
                if len(bullets) > 6:
                    mid = len(bullets) // 2
                    self._pptx_add_bullets(
                        slide,
                        bullets[:mid],
                        Inches(0.8),
                        Inches(2.0),
                        Inches(5.5),
                        TEXT_DARK,
                        PRIMARY,
                    )
                    self._pptx_add_bullets(
                        slide,
                        bullets[mid:],
                        Inches(6.8),
                        Inches(2.0),
                        Inches(5.5),
                        TEXT_DARK,
                        ACCENT2,
                    )
                elif bullets:
                    self._pptx_add_bullets(
                        slide, bullets, Inches(0.8), Inches(2.0), Inches(11.5), TEXT_DARK, PRIMARY
                    )

            # Slide number
            if i > 0:
                nb = slide.shapes.add_textbox(Inches(12.4), Inches(7.0), Inches(0.7), Inches(0.4))
                np = nb.text_frame.paragraphs[0]
                np.text = str(i)
                np.font.size = Pt(10)
                np.font.color.rgb = TEXT_GRAY
                np.alignment = PP_ALIGN.RIGHT

            # Speaker notes — visible only in Presenter view.
            speaker_notes = (slide_data.get("notes") or "").strip()
            if speaker_notes:
                try:
                    slide.notes_slide.notes_text_frame.text = speaker_notes
                except Exception:  # pragma: no cover - defensive
                    pass

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    def _pptx_add_header(self, slide, title_text, slide_width, color, text_color):
        """Add styled header bar to a slide."""
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt

        hdr = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), slide_width, Inches(1.5)
        )
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11.5), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.LEFT

    def _pptx_add_bullets(self, slide, bullets, left, top, width, text_color, dot_color):
        """Add styled bullet list with dot shapes."""
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt

        y = top
        for bullet in bullets:
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, left, y + Inches(0.12), Inches(0.15), Inches(0.15)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = dot_color
            dot.line.fill.background()
            tb = slide.shapes.add_textbox(left + Inches(0.35), y, width - Inches(0.35), Inches(0.6))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(bullet)
            p.font.size = Pt(18)
            p.font.color.rgb = text_color
            y += Inches(0.7)

    def _extract_slides_from_content(self, content: str, title: str) -> list[dict]:
        """Extract slide data. Tries JSON, HTML <section>, markdown, then a text heuristic."""
        import json

        cleaned = (content or "").strip()

        # 1. JSON — either a list of slides or {"slides": [...]}
        try:
            data = json.loads(cleaned)
            if isinstance(data, list) and data:
                return data
            if isinstance(data, dict) and "slides" in data:
                return data["slides"]
        except (json.JSONDecodeError, TypeError):
            pass

        # 2. HTML <section> blocks
        if "<section" in cleaned:
            return self._parse_sections_to_slides(cleaned)

        # 3. Markdown with H1 as deck title and H2 as slide titles (our LLM output shape)
        if _looks_like_markdown_slides(cleaned):
            return _parse_markdown_to_slides(cleaned, fallback_title=title)

        # 4. Last resort — the text heuristic parser
        return self._parse_pptx_content(cleaned, title)

    def _parse_sections_to_slides(self, html: str) -> list[dict]:
        """Parse HTML <section> tags into slide data for PPTX, including tables."""
        slides: list[dict] = []
        # Split by section tags
        sections = re.split(r"<section[^>]*>", html)

        for section in sections:
            if not section.strip():
                continue
            # Remove closing tag
            section = re.sub(r"</section>.*", "", section, flags=re.DOTALL)

            # Extract title from h1 or h2
            title_match = re.search(r"<h[12][^>]*>(.*?)</h[12]>", section, re.DOTALL)
            slide_title = ""
            if title_match:
                slide_title = re.sub(r"<[^>]+>", "", title_match.group(1)).strip()

            # Extract table data if present
            table_data = None
            table_match = re.search(
                r"<table[^>]*>(.*?)</table>", section, re.DOTALL | re.IGNORECASE
            )
            if table_match:
                table_html = table_match.group(1)
                headers: list[str] = []
                rows: list[list[str]] = []
                # Extract headers from <th> tags
                for th in re.finditer(r"<th[^>]*>(.*?)</th>", table_html, re.DOTALL):
                    headers.append(re.sub(r"<[^>]+>", "", th.group(1)).strip())
                # Extract rows from <td> tags grouped by <tr>
                for tr in re.finditer(r"<tr[^>]*>(.*?)</tr>", table_html, re.DOTALL):
                    tr_content = tr.group(1)
                    if "<th" in tr_content:
                        continue  # Skip header row
                    row: list[str] = []
                    for td in re.finditer(r"<td[^>]*>(.*?)</td>", tr_content, re.DOTALL):
                        row.append(re.sub(r"<[^>]+>", "", td.group(1)).strip())
                    if row:
                        rows.append(row)
                if headers and rows:
                    table_data = {"headers": headers, "rows": rows}

            # Extract bullet points from <li> tags
            bullets: list[str] = []
            for li_match in re.finditer(r"<li[^>]*>(.*?)</li>", section, re.DOTALL):
                text = re.sub(r"<[^>]+>", "", li_match.group(1)).strip()
                if text:
                    bullets.append(text)

            # Also extract text from <p>, <div>, <blockquote>, <h3>, <h4> etc.
            # This catches stats, highlights, quotes, and other rich content
            if not bullets or not table_data:
                extra_blocks = re.findall(
                    r"<(?:p|blockquote|h3|h4)[^>]*>(.*?)</(?:p|blockquote|h3|h4)>",
                    section,
                    re.DOTALL,
                )
                for block in extra_blocks:
                    text = re.sub(r"<[^>]+>", "", block).strip()
                    if text and len(text) > 3 and text not in bullets:
                        bullets.append(text)

                # Extract text from div.stat, div.highlight etc.
                div_blocks = re.findall(r"<div[^>]*>(.*?)</div>", section, re.DOTALL)
                for block in div_blocks:
                    # Skip divs that contain other divs (wrappers like .columns)
                    if "<div" in block:
                        # Extract inner text from nested divs
                        inner_texts = re.findall(
                            r"<(?:span|b|p|h3)[^>]*>(.*?)</(?:span|b|p|h3)>", block, re.DOTALL
                        )
                        combined = " ".join(
                            re.sub(r"<[^>]+>", "", t).strip() for t in inner_texts if t.strip()
                        )
                        if combined and len(combined) > 3 and combined not in bullets:
                            bullets.append(combined)
                    else:
                        text = re.sub(r"<[^>]+>", "", block).strip()
                        if text and len(text) > 3 and text not in bullets:
                            bullets.append(text)

            if slide_title or bullets or table_data:
                slide: dict = {"title": slide_title or "Untitled", "bullets": bullets}
                if table_data:
                    slide["table"] = table_data
                slides.append(slide)

        if not slides:
            slides = [{"title": "Presentation", "subtitle": "", "bullets": []}]

        return slides

    def _parse_pptx_content(self, content: str, title: str) -> list[dict]:
        """Parse content into slide data. Supports JSON array or plain text fallback."""
        import json

        # Try JSON first
        try:
            data = json.loads(content)
            if isinstance(data, list) and len(data) > 0:
                return data
            if isinstance(data, dict) and "slides" in data:
                return data["slides"]
        except (json.JSONDecodeError, TypeError):
            pass

        # Fallback: parse HTML/text content into slides by splitting on headings
        slides: list[dict] = []
        current_slide: dict | None = None

        # Strip HTML tags for text extraction
        text = re.sub(r"<[^>]+>", "\n", content)
        text = re.sub(r"\n{3,}", "\n\n", text).strip()
        lines = text.split("\n")

        for line in lines:
            line = line.strip()
            if not line:
                continue

            # Detect heading-like lines (short, no period at end, capitalized)
            is_heading = (
                len(line) < 100
                and not line.endswith(".")
                and (line[0].isupper() or line[0].isdigit())
                and not line.startswith("-")
                and not line.startswith("*")
            )

            if is_heading and (not current_slide or len(current_slide.get("bullets", [])) > 0):
                if current_slide:
                    slides.append(current_slide)
                current_slide = {"title": line, "bullets": []}
            elif current_slide:
                # Clean bullet markers
                clean = re.sub(r"^[-*•]\s*", "", line)
                if clean:
                    current_slide["bullets"].append(clean)
            else:
                current_slide = {"title": title, "bullets": [line]}

        if current_slide:
            slides.append(current_slide)

        # If no slides parsed, create a single title slide
        if not slides:
            slides = [{"title": title, "subtitle": "Generated presentation", "bullets": []}]

        return slides

    def _build_pptx_preview_html(self, title: str, content: str, style: str) -> str:
        """Build an HTML preview for a PPTX presentation from HTML section slides."""
        # Check if content uses <section> tags (new rich HTML format)
        if "<section" in content:
            body = content
        else:
            # Legacy JSON format fallback
            import json

            slides_data = self._parse_pptx_content(content, title)
            slides_parts = []
            for i, slide in enumerate(slides_data):
                slide_title = slide.get("title", f"Slide {i + 1}")
                bullets = slide.get("bullets", [])
                subtitle = slide.get("subtitle", "")
                bullets_html = ""
                if bullets:
                    items = "".join(f"<li>{self._escape_html(b)}</li>" for b in bullets)
                    bullets_html = f"<ul>{items}</ul>"
                subtitle_html = f"<p>{self._escape_html(subtitle)}</p>" if subtitle else ""
                slides_parts.append(
                    f"<section><h2>{self._escape_html(slide_title)}</h2>"
                    f"{subtitle_html}{bullets_html}</section>"
                )
            body = "\n".join(slides_parts)

        return self._wrap_presentation_html(title, body)

    def _wrap_presentation_html(self, title: str, body: str) -> str:
        """Wrap slide sections with reveal.js for interactive slideshow + standalone CSS fallback."""
        css = """
* { box-sizing: border-box; }
.reveal section {
    padding: 40px;
    text-align: left;
}
.reveal h1 {
    font-size: 2.2em;
    font-weight: 800;
    margin-bottom: 0.3em;
    color: #111;
}
.reveal h2 {
    font-size: 1.6em;
    font-weight: 700;
    margin-bottom: 0.5em;
    color: #1a1a1a;
    border-bottom: 3px solid #4f46e5;
    padding-bottom: 8px;
    display: inline-block;
}
.reveal h3 {
    font-size: 1.2em;
    font-weight: 600;
    margin: 0.8em 0 0.4em;
    color: #333;
}
.reveal p {
    font-size: 0.95em;
    margin-bottom: 0.6em;
    color: #444;
    line-height: 1.6;
}
.reveal ul, .reveal ol {
    display: block;
    padding-left: 1.5em;
    margin: 0.5em 0;
}
.reveal li {
    font-size: 0.9em;
    margin-bottom: 0.5em;
    color: #333;
    line-height: 1.5;
}
.reveal table {
    width: 100%;
    border-collapse: collapse;
    margin: 1em 0;
    font-size: 0.85em;
}
.reveal th {
    background: #4f46e5;
    color: white;
    padding: 10px 14px;
    text-align: left;
    font-weight: 600;
}
.reveal td {
    padding: 10px 14px;
    border-bottom: 1px solid #eee;
    color: #333;
}
.reveal tr:nth-child(even) td { background: #f8f8ff; }
.reveal blockquote {
    border-left: 4px solid #4f46e5;
    padding: 12px 16px;
    margin: 1em 0;
    background: #f5f3ff;
    border-radius: 0 8px 8px 0;
    font-style: italic;
    color: #4a4a6a;
    width: 90%;
}
.reveal blockquote cite {
    display: block;
    margin-top: 8px;
    font-style: normal;
    font-size: 0.85em;
    color: #666;
    font-weight: 600;
}
.reveal svg {
    display: block;
    margin: 1em auto;
    max-width: 100%;
}
.columns {
    display: flex;
    gap: 20px;
    margin: 1em 0;
}
.columns > div {
    flex: 1;
    background: #f8fafc;
    padding: 16px;
    border-radius: 10px;
    border: 1px solid #e2e8f0;
}
.columns > div h3 { margin-top: 0; }
.stat {
    display: inline-flex;
    flex-direction: column;
    align-items: center;
    background: #f0f0ff;
    padding: 20px 28px;
    border-radius: 12px;
    margin: 8px 12px 8px 0;
}
.stat .number {
    font-size: 2.5em;
    font-weight: 800;
    color: #4f46e5;
    line-height: 1;
}
.stat .label {
    font-size: 0.75em;
    color: #666;
    margin-top: 6px;
    font-weight: 600;
    text-transform: uppercase;
    letter-spacing: 0.5px;
}
.highlight {
    background: linear-gradient(135deg, #4f46e5 0%, #7c3aed 100%);
    color: white;
    padding: 20px 24px;
    border-radius: 12px;
    margin: 1em 0;
}
.highlight h3, .highlight p, .highlight li { color: white !important; }
.highlight ul { padding-left: 20px; }
.timeline { margin: 1em 0; }
.timeline .event {
    padding: 12px 0 12px 24px;
    border-left: 3px solid #4f46e5;
    margin-bottom: 6px;
    position: relative;
}
.timeline .event::before {
    content: '';
    position: absolute;
    left: -7px;
    top: 16px;
    width: 11px;
    height: 11px;
    border-radius: 50%;
    background: #4f46e5;
}
.timeline .event b { color: #4f46e5; }
.timeline .event p { margin: 4px 0 0; font-size: 0.9em; color: #555; }
"""
        return f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{self._escape_html(title)}</title>
<link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@5.1.0/dist/reveal.css">
<link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@5.1.0/dist/theme/white.css">
<style>{css}</style>
</head>
<body>
<div class="reveal">
<div class="slides">
{body}
</div>
</div>
<script src="https://cdn.jsdelivr.net/npm/reveal.js@5.1.0/dist/reveal.js"></script>
<script>
Reveal.initialize({{
    hash: true,
    slideNumber: true,
    controls: true,
    progress: true,
    center: false,
    transition: 'slide',
    width: 1280,
    height: 720,
}});
</script>
</body>
</html>"""

    async def _upload_bytes(
        self, storage_service: Any, content: bytes, filename: str, path: str
    ) -> dict[str, Any]:
        """Upload raw bytes via the storage client."""
        remote_path = f"{path.strip('/')}/{filename}"
        content_type = CONTENT_TYPES.get(filename.rsplit(".", 1)[-1], "application/octet-stream")
        return await storage_service.upload_bytes(content, remote_path, content_type=content_type)


# Module-level singleton
document_generation_service = DocumentGenerationService()


# ---------------------------------------------------------------------------
# Public orchestration API (used by routes)
# ---------------------------------------------------------------------------


async def create_from_prompt(
    *,
    user_id: str,
    doc_type: str,
    title: str,
    prompt: str,
    format: str = "pdf",
    style: str = "academic",
    course_id: str | None = None,
    topic_id: str | None = None,
):
    """
    Generate a document from a natural-language prompt.

    1. Uses the LLM to generate structured content based on the prompt and document type.
    2. Renders the content to PDF/DOCX/PPTX via ``document_generation_service``.
    3. Persists a ``GeneratedDocument`` record and returns it.

    FREE: PDF only, academic style only.
    PLUS: PDF/DOCX/PPTX, all styles (academic, report, minimal).
    """
    from src.domains.intelligence.reasoning.llm import generate_content
    from src.domains.personal_learning.repository import personal_learning_repo as repo
    from src.domains.personal_learning.services import feature_tier_service, trial_service

    # --- Commercial gate: check format access ---
    if format != "pdf":
        cap_result = await feature_tier_service.check_capability(
            user_id, "document_generation", requested_value=format
        )
        if not cap_result.allowed:
            from fastapi import HTTPException

            raise HTTPException(
                status_code=403,
                detail={
                    "upgradeRequired": True,
                    "reason": cap_result.reason,
                    "capability": cap_result.capability,
                    "upgradeUrl": cap_result.upgrade_url,
                    "trialAvailable": cap_result.trial_available,
                    "upgradeValue": cap_result.upgrade_value,
                },
            )

    # --- Commercial gate: check style access ---
    if style != "academic":
        cap_result = await feature_tier_service.check_capability(
            user_id, "document_generation", requested_value=style
        )
        if not cap_result.allowed:
            from fastapi import HTTPException

            raise HTTPException(
                status_code=403,
                detail={
                    "upgradeRequired": True,
                    "reason": cap_result.reason,
                    "capability": cap_result.capability,
                    "upgradeUrl": cap_result.upgrade_url,
                    "trialAvailable": cap_result.trial_available,
                    "upgradeValue": cap_result.upgrade_value,
                },
            )

    # Record PLUS feature usage if applicable
    if format != "pdf" or style != "academic":
        await trial_service.record_plus_feature_used(user_id, "document_generation")

    llm_prompt = _build_document_prompt(
        doc_type=doc_type, title=title, prompt=prompt, format=format
    )
    max_tokens = _max_tokens_for_type(doc_type)

    try:
        content = await generate_content(llm_prompt, max_tokens=max_tokens, temperature=0.7)
    except Exception as e:
        # Reported as a failure. This used to substitute
        # `# {title}\n\n(Content generation failed. Please try again.)`, render *that* to a PDF,
        # upload it, store it and answer `201` — so a failed generation produced a document in the
        # learner's library whose entire content was an apology, and the client's own copy for this
        # case ("Nothing was saved to your library") was untrue. A 502 rather than a 500: the request
        # was fine and the learner can do nothing differently except try again.
        from fastapi import HTTPException

        logger.error(f"LLM generation failed for document '{title}': {e}")
        raise HTTPException(
            status_code=502,
            detail="The document could not be written. Nothing was saved — please try again.",
        ) from e

    content = _sanitize_llm_output(content)
    if not content.strip():
        from fastapi import HTTPException

        logger.error(f"LLM returned nothing usable for document '{title}'")
        raise HTTPException(
            status_code=502,
            detail="The document came back empty. Nothing was saved — please try again.",
        )

    # Render the file bytes and upload
    result = await document_generation_service.generate_document(
        format=format,
        title=title,
        content=content,
        style=style,
        user_id=user_id,
        doc_type=doc_type,
    )

    # Persist a DB record
    doc = await repo.create_document(
        {
            "userId": user_id,
            "title": title,
            "format": format,
            # The type the learner chose. It shaped the prompt and was then dropped before this
            # dictionary until migration 034, which is why the library page had to guess it back out
            # of the filename.
            "docType": doc_type,
            "style": style,
            "filename": result.get("filename"),
            "fileUrl": result.get("url"),
            "previewUrl": result.get("preview_url"),
            "size": result.get("size", 0),
            "contentType": result.get(
                "content_type", CONTENT_TYPES.get(format, "application/octet-stream")
            ),
            # shareId is always generated but the document stays private until published.
            "shareId": uuid.uuid4().hex[:16],
            "isPublic": False,
        }
    )

    # Record in activity feed
    from src.domains.personal_learning.services import activity_feed_service

    await activity_feed_service.record(
        user_id=user_id,
        activity_type="document_generated",
        title=f"Generated {doc_type}: {title}",
        context={"source": "personal", "docId": doc.id, "format": format},
    )

    # Check milestones
    from src.domains.personal_learning.services import milestone_service

    doc_count = await repo.count_documents_since(user_id, datetime(2000, 1, 1, tzinfo=UTC))
    await milestone_service.check_milestones(user_id, {"documents_generated": doc_count})

    return doc


async def get_document(*, user_id: str, doc_id: str):
    """Fetch a generated document by id."""
    from src.domains.personal_learning.repository import personal_learning_repo as repo
    from src.shared.exceptions import NotFoundError

    doc = await repo.find_document(doc_id, user_id)
    if not doc:
        raise NotFoundError("Document", doc_id)
    return doc


async def publish_document(*, user_id: str, doc_id: str):
    """Mark a document as public and return a share id.

    The share id is rotated on every publish. A document published, unpublished and published again
    therefore gets a *new* link rather than reviving the old one, which is what makes unpublishing
    mean something.
    """
    import uuid as _uuid

    from src.domains.personal_learning.repository import personal_learning_repo as repo
    from src.shared.exceptions import NotFoundError

    doc = await repo.find_document(doc_id, user_id)
    if not doc:
        raise NotFoundError("Document", doc_id)

    share_id = _uuid.uuid4().hex[:16]
    await repo.update_document(doc_id, {"isPublic": True, "shareId": share_id})
    return await repo.find_document(doc_id, user_id)


async def unpublish_document(*, user_id: str, doc_id: str):
    """Withdraw a published document and retire the link it was shared under.

    ``isPublic`` goes false and the share id is rotated, so the URL the learner sent out stops
    resolving permanently — a later republish issues a different one. ``shareId`` is `NOT NULL`, so
    it is replaced rather than cleared; a private document has always had one.

    What this does **not** do, stated because the alternative is a button that claims more than it
    delivers: it does not revoke the file. Documents are stored at unauthenticated public URLs, so
    anybody already holding ``fileUrl`` or ``previewUrl`` keeps their access, and anybody who
    downloaded the file keeps the file. Unpublishing withdraws the shared *page*. Making the file
    itself revocable needs token-authenticated storage URLs or a proxied download route, which is a
    change to how every stored object is served, not to this function.

    Idempotent: unpublishing something already private is a success, and still rotates the id, since
    the caller's intent is "this link should not work".
    """
    import uuid as _uuid

    from src.domains.personal_learning.repository import personal_learning_repo as repo
    from src.shared.exceptions import NotFoundError

    doc = await repo.find_document(doc_id, user_id)
    if not doc:
        raise NotFoundError("Document", doc_id)

    await repo.update_document(doc_id, {"isPublic": False, "shareId": _uuid.uuid4().hex[:16]})
    return await repo.find_document(doc_id, user_id)


async def delete_document(*, user_id: str, doc_id: str) -> None:
    """Delete a document, and the two objects it stored.

    **Hard delete, not a `deletedAt` column.** Every generated document is two stored objects — the
    file and its HTML preview — so a soft delete keeps paying for storage on something the learner
    has said they do not want, forever, and hides it behind a filter that every read path then has to
    remember. The learner's other artifacts here already delete outright, notes included, and a
    document is reproducible: the prompt that made it is the thing worth keeping, and it is in the
    activity feed.

    Storage first, row second, following ``study_plan_service.delete_material``. A row pointing at a
    file that still exists is recoverable — the learner sees the document and can try again. A
    deleted row pointing at a live object leaves something nobody can find, name or clean up. Both
    objects are attempted even if the first fails, so a transient failure on the preview does not
    strand the file.
    """
    from src.domains.personal_learning.repository import personal_learning_repo as repo
    from src.shared.exceptions import NotFoundError
    from src.shared.infrastructure.storage import storage_service

    doc = await repo.find_document(doc_id, user_id)
    if not doc:
        raise NotFoundError("Document", doc_id)

    for url in (doc.file_url, doc.preview_url):
        if not url or not storage_service.owns_url(url):
            continue
        removed = await storage_service.delete(url)
        if not removed:
            # Logged rather than raised. Storage being unreachable should not leave the learner
            # unable to remove a document from their own library, and the row is the record that
            # makes the object findable — keeping it because a delete failed inverts the problem.
            logger.warning(
                "Document object could not be deleted from storage; deleting the row anyway",
                extra={"user_id": user_id, "doc_id": doc_id, "url": url},
            )

    await repo.delete_document(doc_id, user_id)


async def list_documents(
    *,
    user_id: str,
    page: int = 1,
    page_size: int = 20,
    search: str | None = None,
    format: str | None = None,
    type: str | None = None,
):
    """Return a page of generated documents for the current user, newest first.

    ``search`` matches title or filename; ``format`` and ``type`` are exact. All three are applied in
    the query, not after the fact: the page used to be filtered in the browser over whichever twenty
    documents it had loaded, so searching a library of sixty found nothing in the other forty — and
    the type it filtered on was inferred from the filename rather than stored at all.
    """
    from src.domains.personal_learning.repository import personal_learning_repo as repo

    skip = max(0, (page - 1) * page_size)
    return await repo.list_documents(
        user_id,
        skip=skip,
        take=page_size,
        search=search,
        doc_format=format,
        doc_type=type,
    )


async def get_summary(*, user_id: str) -> dict[str, Any]:
    """Library-wide document figures: totals, published, this month, and the format breakdown.

    "This month" is resolved in the learner's own timezone where we have one. ``UserPreferences.timezone``
    is `NOT NULL` with a `"UTC"` default, so it is only a fact when the source says it was observed —
    ``resolve_learner_timezone`` is what encodes that, and it falls back to UTC flagged as unknown. The
    boundary is returned alongside the count so the page can say what it measured from rather than
    asserting a calendar month for a learner whose location we have never been told.
    """
    import asyncio
    from datetime import UTC as _UTC
    from datetime import datetime as _datetime

    from src.domains.personal_learning.repository import personal_learning_repo as repo
    from src.shared.time.learner_timezone import resolve_learner_timezone, to_learner_local

    learner_timezone = await resolve_learner_timezone(user_id)
    local_now = to_learner_local(_datetime.now(_UTC), learner_timezone)
    local_month_start = local_now.replace(day=1, hour=0, minute=0, second=0, microsecond=0)
    month_start = local_month_start.astimezone(_UTC)

    counts, formats = await asyncio.gather(
        repo.count_documents(user_id, since=month_start),
        repo.list_document_formats(user_id),
    )
    total, published, created_this_month = counts

    return {
        "total": total,
        "published": published,
        "createdThisMonth": created_this_month,
        "monthStart": month_start,
        "formats": [{"format": fmt, "count": count} for fmt, count in formats],
    }


async def get_by_share_id(*, share_id: str, requester_id: str | None = None):
    """
    Fetch a document by its share id.

    - Public documents (``is_public=True``) are visible to anyone.
    - Private documents are only returned when ``requester_id`` matches the owner,
      so an owner can preview their share URL before publishing.
    - Returns ``None`` if the document does not exist or the caller has no access.
    """
    from src.domains.personal_learning.repository import personal_learning_repo as repo

    doc = await repo.find_document_by_share_id(share_id)
    if not doc:
        return None
    if getattr(doc, "is_public", False):
        return doc
    if requester_id and getattr(doc, "user_id", None) == requester_id:
        return doc
    return None


# Backwards-compatible alias for older callers.
async def get_shared_document(share_id: str):
    return await get_by_share_id(share_id=share_id)


# ---------------------------------------------------------------------------
# Prompt engineering per document type
# ---------------------------------------------------------------------------

_DOC_TYPE_GUIDANCE: dict[str, dict[str, str]] = {
    "essay": {
        "word_target": "800 to 1200 words",
        "voice": "formal and analytical",
        "structure": (
            "an introduction with a clear thesis, three to five body paragraphs each "
            "developing one argument with evidence, and a conclusion that ties the thesis "
            "to a broader implication."
        ),
    },
    "report": {
        "word_target": "1200 to 2000 words",
        "voice": "objective, precise, and professional",
        "structure": (
            "an executive summary, introduction, methodology, findings grouped under "
            "sub-headings, discussion, actionable recommendations, and a brief conclusion."
        ),
    },
    "presentation": {
        "word_target": "8 to 14 slides",
        "voice": "concise and speaker-friendly",
        "structure": (
            "one `# Deck Title` line and an optional subtitle paragraph, then one slide "
            "per `## Slide Title`. Give each slide 3 to 6 short bullets. To make the deck "
            "visually varied, tag some slides with a layout using one of these prefixes "
            "in the H2 title:\n"
            "  - `## [section] Section Name` for a full-bleed section divider\n"
            "  - `## [big] 87%` with the caption as a single bullet, for a big-number stat\n"
            '  - `## [quote] "Design is intelligence made visible."` with the author as the only bullet\n'
            "  - `## [two-col | Pros | Cons] Comparison title` with the bullets split evenly between the two columns\n"
            "Use these tags sparingly (2 to 4 across the deck) to add rhythm. End with a "
            "`## Summary` slide and a `## [closer] Thank you` slide."
        ),
    },
    "letter": {
        "word_target": "150 to 400 words",
        "voice": "polite and direct",
        "structure": (
            "sender contact, date, recipient contact, salutation, one to three body "
            "paragraphs stating the purpose and any request, a closing line, and a sign-off."
        ),
    },
    "cv": {
        "word_target": "one page, dense",
        "voice": "confident and results-focused",
        "structure": (
            "contact block, a short professional summary, reverse-chronological experience "
            "with three to five bullets per role starting with action verbs and including "
            "measurable results, education, and grouped skills. Use bracketed placeholders "
            "like [Company] for any facts the user did not supply."
        ),
    },
}


_DEFAULT_GUIDANCE = {
    "word_target": "600 to 1000 words",
    "voice": "clear and professional",
    "structure": "clear headings for major sections and short focused paragraphs.",
}


def _resolve_doc_type_guidance(doc_type: str) -> dict[str, str]:
    """Return the writing guidance block for a doc type (case-insensitive)."""
    return _DOC_TYPE_GUIDANCE.get((doc_type or "").strip().lower(), _DEFAULT_GUIDANCE)


def _build_document_prompt(*, doc_type: str, title: str, prompt: str, format: str) -> str:
    """Compose a short, focused LLM prompt for document generation.

    Kept intentionally lean — long rule-lists cause the model to echo
    instructions back as content. Voice is enforced via the sanitizer
    (see ``_sanitize_llm_output``), not the prompt.
    """
    guidance = _resolve_doc_type_guidance(doc_type)

    article = "an" if (doc_type or "").strip().lower()[:1] in "aeiou" else "a"
    return (
        f'Write {article} {doc_type} titled "{title}".\n\n'
        f"{prompt.strip()}\n\n"
        f"Length: {guidance['word_target']}. Voice: {guidance['voice']}.\n"
        f"Structure it as: {guidance['structure']}\n\n"
        f"Return only the finished document in Markdown. "
        f"Use plain text for math (e.g. O(log n), not LaTeX). "
        f"Use commas and periods, not em-dashes."
    )


def _max_tokens_for_type(doc_type: str) -> int:
    """Budget output tokens based on the expected length of the doc type."""
    key = (doc_type or "").strip().lower()
    return {
        "letter": 1500,
        "cv": 2500,
        "essay": 3000,
        "presentation": 3500,
        "report": 5000,
    }.get(key, 3000)


# ---------------------------------------------------------------------------
# LLM output sanitization
# ---------------------------------------------------------------------------

# Compiled once at module load.
_EMDASH_PATTERN = re.compile(r"\s*[—–]\s*")

# LaTeX inline math: $...$ but not $$...$$ (display) — we strip both.
_LATEX_INLINE_PATTERN = re.compile(r"\$([^$\n]+?)\$")
_LATEX_DISPLAY_PATTERN = re.compile(r"\$\$([^$]+?)\$\$", re.DOTALL)
_LATEX_PAREN_PATTERN = re.compile(r"\\\(([^)]+?)\\\)")
_LATEX_BRACKET_PATTERN = re.compile(r"\\\[([^\]]+?)\\\]", re.DOTALL)

# Common LaTeX macros we strip when we hit plain-text math.
_LATEX_MACRO_REPLACEMENTS = {
    r"\log": "log",
    r"\ln": "ln",
    r"\sin": "sin",
    r"\cos": "cos",
    r"\tan": "tan",
    r"\infty": "infinity",
    r"\times": "×",
    r"\cdot": "·",
    r"\pm": "±",
    r"\leq": "≤",
    r"\geq": "≥",
    r"\neq": "≠",
    r"\approx": "≈",
    r"\alpha": "alpha",
    r"\beta": "beta",
    r"\gamma": "gamma",
    r"\delta": "delta",
    r"\theta": "theta",
    r"\pi": "pi",
    r"\sigma": "sigma",
    r"\mu": "mu",
    r"\lambda": "lambda",
}


def _strip_latex(text: str) -> str:
    """Convert LaTeX math to plain text since our HTML renderer does not process MathJax/KaTeX."""

    def _unwrap(match: re.Match) -> str:
        inner = match.group(1).strip()
        for macro, replacement in _LATEX_MACRO_REPLACEMENTS.items():
            inner = inner.replace(macro, replacement)
        # Strip leftover backslashes on unknown macros.
        inner = re.sub(r"\\[a-zA-Z]+\s*", "", inner)
        # Normalize whitespace.
        inner = re.sub(r"\s+", " ", inner).strip()
        return inner

    text = _LATEX_DISPLAY_PATTERN.sub(_unwrap, text)
    text = _LATEX_INLINE_PATTERN.sub(_unwrap, text)
    text = _LATEX_BRACKET_PATTERN.sub(_unwrap, text)
    text = _LATEX_PAREN_PATTERN.sub(_unwrap, text)
    return text


def _sanitize_llm_output(text: str) -> str:
    """
    Post-process LLM output to remove the most obvious AI-writing tells.

    - Replaces em-dashes / en-dashes with a comma.
    - Strips LaTeX math notation ($...$, $$...$$, \\(...\\), \\[...\\]).
    - Removes leaked markdown code fences around the document.
    """
    if not text:
        return text

    cleaned = text.strip()

    # Drop leading/trailing markdown fences (```markdown ... ```).
    if cleaned.startswith("```"):
        lines = cleaned.split("\n", 1)
        cleaned = lines[1] if len(lines) > 1 else ""
    if cleaned.endswith("```"):
        cleaned = cleaned.rsplit("```", 1)[0]

    cleaned = _strip_latex(cleaned)
    cleaned = _EMDASH_PATTERN.sub(", ", cleaned)

    return cleaned.strip()


# ---------------------------------------------------------------------------
# Format normalization
# ---------------------------------------------------------------------------

_FORMAT_ALIASES: dict[str, str] = {
    "pdf": "pdf",
    "docx": "docx",
    "doc": "docx",
    "word": "docx",
    # PowerPoint output is disabled for now — the HTML/PDF slide renderer
    # produces much better looking decks. All slide requests fall back to PDF.
    "pptx": "pdf",
    "ppt": "pdf",
    "powerpoint": "pdf",
    "slides": "pdf",
    "presentation": "pdf",
}


def _normalize_format(value: str | None) -> str:
    """Map user-supplied format strings to canonical values (pdf, docx, pptx)."""
    key = (value or "").strip().lower()
    if key not in _FORMAT_ALIASES:
        raise ValueError(f"Unsupported format: {value}. Use 'pdf', 'docx', or 'pptx'.")
    return _FORMAT_ALIASES[key]


# ---------------------------------------------------------------------------
# Markdown → slide parser (for PPTX)
# ---------------------------------------------------------------------------

_MD_H1 = re.compile(r"^#\s+(.+)$")
_MD_H2 = re.compile(r"^##\s+(.+)$")
_MD_H3 = re.compile(r"^###\s+(.+)$")
_MD_BULLET = re.compile(r"^\s*[-*+]\s+(.+)$")
_MD_NUMBERED = re.compile(r"^\s*\d+\.\s+(.+)$")
_MD_SPEAKER_NOTES = re.compile(r"^\s*(?:speaker\s+notes?|notes?)\s*[:：]\s*(.+)$", re.IGNORECASE)


def _looks_like_markdown_slides(content: str) -> bool:
    """Heuristic: content is markdown-style slides if it has at least one H2 heading."""
    for line in content.splitlines():
        if _MD_H2.match(line.strip()):
            return True
    return False


def _parse_markdown_to_slides(content: str, *, fallback_title: str) -> list[dict]:
    """
    Parse a markdown document with H1 deck title and H2 per slide into slide dicts.

    Recognized structure:

        # Deck Title                       -> title slide
        Optional subtitle line
        ## Slide 1 title                   -> new slide
        - bullet one
        - bullet two
        Speaker notes: extra context...    -> attached to previous slide as `notes`
        ## Slide 2 title
        ...

    Returns a list of ``{title, bullets, subtitle?, notes?}`` dicts. Always
    includes a title slide as the first element.
    """
    lines = content.splitlines()
    deck_title: str | None = None
    subtitle_lines: list[str] = []
    slides: list[dict] = []
    current: dict | None = None
    in_preamble = True  # everything before the first H2 is deck-level

    def _push(slide: dict | None) -> None:
        if not slide:
            return
        # Drop empty slides
        if not slide.get("title") and not slide.get("bullets"):
            return
        slides.append(slide)

    for raw in lines:
        line = raw.rstrip()
        stripped = line.strip()
        if not stripped:
            continue

        # H1 → deck title
        m = _MD_H1.match(stripped)
        if m:
            deck_title = _clean_inline(m.group(1))
            continue

        # H2 → new slide
        m = _MD_H2.match(stripped)
        if m:
            in_preamble = False
            _push(current)
            current = {"title": _clean_inline(m.group(1)), "bullets": [], "notes": ""}
            continue

        # H3 → treat as bold callout inside current slide (also opens a slide if none)
        m = _MD_H3.match(stripped)
        if m:
            in_preamble = False
            if current is None:
                current = {"title": _clean_inline(m.group(1)), "bullets": [], "notes": ""}
                continue
            # Attach as a sub-item
            current["bullets"].append(_clean_inline(m.group(1)))
            continue

        # Speaker notes → attach to the current slide
        m = _MD_SPEAKER_NOTES.match(stripped)
        if m and current is not None:
            note = _clean_inline(m.group(1))
            current["notes"] = (current["notes"] + " " + note).strip() if current["notes"] else note
            continue

        # Bullet (- or *)
        m = _MD_BULLET.match(stripped) or _MD_NUMBERED.match(stripped)
        if m:
            in_preamble = False
            if current is None:
                # Bullets before any H2 — attach to a synthetic first slide
                current = {"title": deck_title or fallback_title, "bullets": [], "notes": ""}
            current["bullets"].append(_clean_inline(m.group(1)))
            continue

        # Plain paragraph
        if in_preamble:
            subtitle_lines.append(_clean_inline(stripped))
        elif current is not None:
            # A prose line inside a slide → treat as a bullet if we do not already have many,
            # otherwise fold into speaker notes.
            text = _clean_inline(stripped)
            if len(current["bullets"]) < 6:
                current["bullets"].append(text)
            else:
                current["notes"] = (
                    (current["notes"] + " " + text).strip() if current["notes"] else text
                )

    _push(current)

    # Build the title slide, then the body slides.
    result: list[dict] = [
        {
            "title": deck_title or fallback_title,
            "subtitle": " ".join(subtitle_lines).strip() or None,
            "bullets": [],
            "notes": "",
        }
    ]
    result.extend(slides)

    # Ensure a closing slide exists.
    def _looks_like_closer(t: str) -> bool:
        lower = (t or "").lower()
        # Strip a leading layout tag like "[closer] "
        lower = re.sub(r"^\s*\[[^\]]+\]\s*", "", lower).strip()
        return lower in {"thank you", "thanks", "questions", "q&a", "questions?"}

    if not result or not _looks_like_closer(result[-1].get("title", "")):
        result.append(
            {
                "title": "Thank you",
                "subtitle": "Questions?",
                "bullets": [],
                "notes": "",
            }
        )

    return result


_INLINE_MD_BOLD = re.compile(r"\*\*(.+?)\*\*")
_INLINE_MD_ITALIC = re.compile(r"(?<!\*)\*([^*\n]+?)\*(?!\*)")
_INLINE_MD_CODE = re.compile(r"`([^`\n]+?)`")


def _clean_inline(text: str) -> str:
    """Strip inline markdown markers so the raw text renders cleanly in slides."""
    if not text:
        return ""
    text = _INLINE_MD_BOLD.sub(r"\1", text)
    text = _INLINE_MD_ITALIC.sub(r"\1", text)
    text = _INLINE_MD_CODE.sub(r"\1", text)
    return text.strip()


# ---------------------------------------------------------------------------
# HTML slide renderer for beautiful PDF presentations
# ---------------------------------------------------------------------------

# CSS palette — kept in one place so themes are easy to swap.
_SLIDE_PALETTES = {
    "indigo": {
        "bg": "#0F172A",  # slide title/section bg
        "bg_light": "#FFFFFF",  # content slide bg
        "accent": "#6366F1",  # primary accent
        "accent2": "#22D3EE",  # secondary accent
        "text_dark": "#0F172A",
        "text_muted": "#64748B",
        "text_light": "#F8FAFC",
        "border": "#E2E8F0",
    },
    "emerald": {
        "bg": "#022C22",
        "bg_light": "#FFFFFF",
        "accent": "#10B981",
        "accent2": "#F59E0B",
        "text_dark": "#022C22",
        "text_muted": "#64748B",
        "text_light": "#ECFDF5",
        "border": "#E2E8F0",
    },
    "rose": {
        "bg": "#1F0A15",
        "bg_light": "#FFFFFF",
        "accent": "#F43F5E",
        "accent2": "#F59E0B",
        "text_dark": "#1F0A15",
        "text_muted": "#64748B",
        "text_light": "#FFF1F2",
        "border": "#E2E8F0",
    },
}


_SLIDE_CSS_TEMPLATE = """
@page {
    size: 33.87cm 19.05cm; /* 16:9 landscape (widescreen) */
    margin: 0;
}
* { box-sizing: border-box; margin: 0; padding: 0; }
body {
    font-family: -apple-system, BlinkMacSystemFont, "SF Pro Display", "Inter", "Helvetica Neue", Helvetica, Arial, sans-serif;
    background: {bg_light};
    color: {text_dark};
    -webkit-font-smoothing: antialiased;
}
.slide {
    width: 33.87cm;
    height: 19.05cm;
    padding: 2.2cm 2.8cm;
    page-break-after: always;
    page-break-inside: avoid;
    position: relative;
    display: flex;
    flex-direction: column;
    justify-content: flex-start;
    overflow: hidden;
}
.slide:last-child { page-break-after: auto; }

/* Slide numbers and brand */
.slide-num {
    position: absolute;
    right: 1.4cm;
    bottom: 1.0cm;
    font-size: 10pt;
    color: {text_muted};
    letter-spacing: 0.1em;
}
.brand {
    position: absolute;
    left: 1.4cm;
    bottom: 1.0cm;
    font-size: 10pt;
    color: {text_muted};
    letter-spacing: 0.15em;
    text-transform: uppercase;
    font-weight: 600;
}

/* Slide types */
.slide-title {
    background: {bg};
    color: {text_light};
    justify-content: center;
    align-items: flex-start;
    padding-left: 3.6cm;
    padding-right: 3.6cm;
}
.slide-title h1 {
    font-size: 60pt;
    font-weight: 800;
    line-height: 1.05;
    max-width: 22cm;
    letter-spacing: -0.02em;
}
.slide-title .subtitle {
    font-size: 22pt;
    color: rgba(255, 255, 255, 0.7);
    margin-top: 0.6cm;
    max-width: 22cm;
    font-weight: 400;
    line-height: 1.35;
}
.slide-title .accent-bar {
    width: 6cm;
    height: 4px;
    background: {accent};
    margin-bottom: 1.2cm;
    border-radius: 2px;
}
.slide-title .brand,
.slide-title .slide-num { color: rgba(255, 255, 255, 0.5); }

.slide-section {
    background: {accent};
    color: {text_light};
    justify-content: center;
    align-items: flex-start;
}
.slide-section .eyebrow {
    font-size: 12pt;
    letter-spacing: 0.2em;
    text-transform: uppercase;
    color: rgba(255, 255, 255, 0.75);
    margin-bottom: 0.6cm;
    font-weight: 600;
}
.slide-section h2 {
    font-size: 56pt;
    font-weight: 800;
    line-height: 1.05;
    letter-spacing: -0.02em;
    max-width: 24cm;
}
.slide-section .brand,
.slide-section .slide-num { color: rgba(255, 255, 255, 0.6); }

.slide-content { }
.slide-content .eyebrow {
    font-size: 11pt;
    letter-spacing: 0.2em;
    text-transform: uppercase;
    color: {accent};
    margin-bottom: 0.35cm;
    font-weight: 700;
}
.slide-content h2 {
    font-size: 34pt;
    font-weight: 700;
    line-height: 1.15;
    letter-spacing: -0.01em;
    margin-bottom: 1.0cm;
    max-width: 24cm;
}
.slide-content ul {
    list-style: none;
    padding: 0;
    margin: 0;
    max-width: 24cm;
}
.slide-content li {
    font-size: 18pt;
    line-height: 1.5;
    padding: 0.32cm 0 0.32cm 1.1cm;
    position: relative;
    color: {text_dark};
    border-bottom: 1px solid {border};
}
.slide-content li:last-child { border-bottom: none; }
.slide-content li::before {
    content: "";
    position: absolute;
    left: 0;
    top: 0.75cm;
    width: 0.55cm;
    height: 2px;
    background: {accent};
    border-radius: 1px;
}

.slide-quote {
    justify-content: center;
    align-items: flex-start;
    background: {bg_light};
}
.slide-quote .mark {
    font-size: 96pt;
    line-height: 0.6;
    color: {accent};
    font-family: Georgia, serif;
    margin-bottom: 0.3cm;
}
.slide-quote blockquote {
    font-family: Georgia, "Times New Roman", serif;
    font-size: 32pt;
    line-height: 1.35;
    color: {text_dark};
    max-width: 26cm;
    font-weight: 400;
    font-style: italic;
}
.slide-quote .attribution {
    margin-top: 0.9cm;
    font-size: 14pt;
    color: {text_muted};
    letter-spacing: 0.05em;
}

.slide-big {
    justify-content: center;
    align-items: flex-start;
}
.slide-big .number {
    font-size: 110pt;
    font-weight: 900;
    line-height: 1.0;
    color: {accent};
    letter-spacing: -0.04em;
    max-width: 28cm;
    /* Prevent awkward mid-token wrapping like "O(log n)" splitting to two lines */
    white-space: nowrap;
    overflow-wrap: normal;
}
/* Auto-shrink for anything longer than a short stat by dropping to a smaller base */
.slide-big .number.long {
    font-size: 68pt;
    white-space: normal;
    line-height: 1.05;
}
.slide-big .caption {
    font-size: 22pt;
    color: {text_dark};
    margin-top: 0.8cm;
    max-width: 26cm;
    line-height: 1.3;
    font-weight: 500;
}
.slide-big .eyebrow {
    font-size: 11pt;
    letter-spacing: 0.2em;
    text-transform: uppercase;
    color: {text_muted};
    margin-bottom: 0.4cm;
    font-weight: 700;
}

.slide-two-col { }
.slide-two-col h2 {
    font-size: 32pt;
    font-weight: 700;
    line-height: 1.15;
    letter-spacing: -0.01em;
    margin-bottom: 0.9cm;
    max-width: 28cm;
}
.slide-two-col .cols {
    display: table;
    width: 100%;
    border-spacing: 0.9cm 0;
}
.slide-two-col .col {
    display: table-cell;
    vertical-align: top;
    width: 50%;
}
.slide-two-col .col-heading {
    font-size: 14pt;
    letter-spacing: 0.15em;
    text-transform: uppercase;
    color: {accent};
    font-weight: 700;
    margin-bottom: 0.4cm;
    padding-bottom: 0.25cm;
    border-bottom: 2px solid {accent};
}
.slide-two-col .col ul {
    list-style: none;
    padding: 0;
    margin: 0;
}
.slide-two-col .col li {
    font-size: 15pt;
    line-height: 1.5;
    padding: 0.2cm 0 0.2cm 0.7cm;
    position: relative;
    color: {text_dark};
}
.slide-two-col .col li::before {
    content: "";
    position: absolute;
    left: 0;
    top: 0.6cm;
    width: 0.35cm;
    height: 2px;
    background: {accent};
    border-radius: 1px;
}

.slide-closer {
    background: {bg};
    color: {text_light};
    justify-content: center;
    align-items: center;
    text-align: center;
}
.slide-closer h2 {
    font-size: 72pt;
    font-weight: 800;
    letter-spacing: -0.02em;
    line-height: 1.05;
}
.slide-closer .subtitle {
    font-size: 22pt;
    color: rgba(255, 255, 255, 0.7);
    margin-top: 0.8cm;
    max-width: 24cm;
}
.slide-closer .brand,
.slide-closer .slide-num { color: rgba(255, 255, 255, 0.5); }
"""


# Regex to detect a layout tag prefix on an H2, e.g. `## [big] 87%`
_LAYOUT_TAG_PATTERN = re.compile(
    r"^\[\s*([a-z0-9_-]+)(?:\s*\|\s*([^\]]+))?\s*\]\s*(.*)$", re.IGNORECASE
)


def _apply_layout_tags(slides: list[dict]) -> list[dict]:
    """
    Look at each slide title for a leading `[layout]` or `[layout | arg | arg]` tag
    and attach layout metadata. Removes the tag from the visible title.
    """
    for slide in slides:
        title = (slide.get("title") or "").strip()
        m = _LAYOUT_TAG_PATTERN.match(title)
        if not m:
            slide.setdefault("layout", "content")
            continue
        layout = m.group(1).lower()
        args_raw = m.group(2) or ""
        stripped_title = m.group(3).strip()
        slide["layout"] = layout
        if args_raw.strip():
            slide["layout_args"] = [a.strip() for a in args_raw.split("|") if a.strip()]
        slide["title"] = stripped_title
    return slides


def _render_slide_html(slide: dict, index: int, total: int, deck_title: str) -> str:
    """Render a single slide dict into an HTML `<section class="slide ...">` block."""
    layout = (slide.get("layout") or "content").lower()
    title = _html_escape(slide.get("title") or "")
    subtitle = _html_escape(slide.get("subtitle") or "")
    bullets = slide.get("bullets") or []
    args = slide.get("layout_args") or []

    slide_num_html = (
        f'<span class="slide-num">{index:02d} / {total:02d}</span>' if index > 0 else ""
    )
    brand_html = f'<span class="brand">{_html_escape(deck_title)}</span>' if index > 0 else ""

    # First slide is always the deck title.
    if index == 0:
        subtitle_html = f'<div class="subtitle">{subtitle}</div>' if subtitle else ""
        return (
            f'<section class="slide slide-title">'
            f'<div class="accent-bar"></div>'
            f"<h1>{title}</h1>{subtitle_html}"
            f"{brand_html}{slide_num_html}"
            f"</section>"
        )

    # Closing slides
    if layout == "closer" or title.lower() in {"thank you", "questions", "q&a"}:
        subtitle_html = f'<div class="subtitle">{subtitle or "Any questions?"}</div>'
        return (
            f'<section class="slide slide-closer">'
            f'<h2>{title or "Thank you"}</h2>{subtitle_html}'
            f"{brand_html}{slide_num_html}"
            f"</section>"
        )

    if layout in {"section", "divider"}:
        eyebrow = f'<div class="eyebrow">Section {_section_number(index)}</div>'
        return (
            f'<section class="slide slide-section">{eyebrow}<h2>{title}</h2>'
            f"{brand_html}{slide_num_html}</section>"
        )

    if layout == "quote":
        # bullets[0] is the attribution when present
        attribution = _html_escape(bullets[0]) if bullets else ""
        attribution_html = f'<div class="attribution">— {attribution}</div>' if attribution else ""
        return (
            f'<section class="slide slide-quote">'
            f'<div class="mark">&ldquo;</div>'
            f"<blockquote>{title}</blockquote>{attribution_html}"
            f"{brand_html}{slide_num_html}</section>"
        )

    if layout in {"big", "big-number", "stat"}:
        eyebrow = f'<div class="eyebrow">{_html_escape(args[0])}</div>' if args else ""
        caption = " ".join(_html_escape(b) for b in bullets)
        caption_html = f'<div class="caption">{caption}</div>' if caption else ""
        # Short glyphs get the huge treatment; anything longer scales down.
        number_class = "number" if len(title) <= 8 else "number long"
        return (
            f'<section class="slide slide-big">{eyebrow}'
            f'<div class="{number_class}">{title}</div>{caption_html}'
            f"{brand_html}{slide_num_html}</section>"
        )

    if layout in {"two-col", "twocol", "compare", "comparison"}:
        left_title = _html_escape(args[0]) if len(args) >= 1 else "Left"
        right_title = _html_escape(args[1]) if len(args) >= 2 else "Right"
        # Split bullets into two halves; caller can also send a full-list slide.
        half = (len(bullets) + 1) // 2
        left_bullets = bullets[:half]
        right_bullets = bullets[half:]
        left_html = "".join(f"<li>{_html_escape(b)}</li>" for b in left_bullets)
        right_html = "".join(f"<li>{_html_escape(b)}</li>" for b in right_bullets)
        return (
            f'<section class="slide slide-two-col">'
            f"<h2>{title}</h2>"
            f'<div class="cols">'
            f'<div class="col"><div class="col-heading">{left_title}</div><ul>{left_html}</ul></div>'
            f'<div class="col"><div class="col-heading">{right_title}</div><ul>{right_html}</ul></div>'
            f"</div>{brand_html}{slide_num_html}</section>"
        )

    # Default: content slide with bullets.
    eyebrow = f'<div class="eyebrow">{_section_number(index)}</div>'
    items_html = "".join(f"<li>{_html_escape(b)}</li>" for b in bullets)
    return (
        f'<section class="slide slide-content">{eyebrow}<h2>{title}</h2>'
        f"<ul>{items_html}</ul>"
        f"{brand_html}{slide_num_html}</section>"
    )


def _section_number(index: int) -> str:
    return f"{index:02d}"


def _html_escape(value: str) -> str:
    if value is None:
        return ""
    return (
        str(value)
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def build_presentation_html(deck_title: str, slides: list[dict], theme: str = "indigo") -> str:
    """Compose the full HTML document for a themed slide deck."""
    palette = _SLIDE_PALETTES.get(theme, _SLIDE_PALETTES["indigo"])
    css = _SLIDE_CSS_TEMPLATE
    for key, value in palette.items():
        css = css.replace("{" + key + "}", value)
    slides = _apply_layout_tags(slides)
    total = len(slides)
    body = "\n".join(
        _render_slide_html(slide, i, total, deck_title) for i, slide in enumerate(slides)
    )
    return (
        "<!DOCTYPE html>"
        f'<html lang="en"><head><meta charset="utf-8">'
        f"<title>{_html_escape(deck_title)}</title>"
        f"<style>{css}</style></head>"
        f"<body>{body}</body></html>"
    )


def render_presentation_pdf(deck_title: str, slides: list[dict], theme: str = "indigo") -> bytes:
    """Render a slide deck to a widescreen PDF using WeasyPrint."""
    html = build_presentation_html(deck_title, slides, theme=theme)
    try:
        from weasyprint import HTML as _WeasyHTML

        return _WeasyHTML(string=html).write_pdf()
    except (ImportError, OSError) as e:  # pragma: no cover - depends on system libs
        logger.warning(
            f"WeasyPrint unavailable ({type(e).__name__}: {e}). "
            "Falling back to xhtml2pdf for slides."
        )
        from xhtml2pdf import pisa

        buf = io.BytesIO()
        pisa.CreatePDF(io.StringIO(html), dest=buf)
        buf.seek(0)
        return buf.getvalue()
