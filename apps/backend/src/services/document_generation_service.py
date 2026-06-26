"""
Document Generation Service.

Generates downloadable documents (PDF, DOCX) from HTML content.
Uses WeasyPrint for PDF rendering and htmldocx for DOCX conversion.
"""

from __future__ import annotations

import io
import logging
import re
import uuid
from datetime import UTC, datetime
from typing import Any

logger = logging.getLogger(__name__)

# Content type mappings
CONTENT_TYPES = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

# Max content length to prevent abuse (100k chars ~= 50 pages)
MAX_CONTENT_LENGTH = 100_000

# CSS styles for different document styles
_ACADEMIC_CSS = """
@page {
    size: A4;
    margin: 2.5cm 2cm;
}
body {
    font-family: 'Times New Roman', Times, Georgia, serif;
    font-size: 12pt;
    line-height: 1.6;
    color: #1a1a1a;
}
h1 { font-size: 22pt; margin-top: 1.5em; margin-bottom: 0.5em; color: #111; }
h2 { font-size: 16pt; margin-top: 1.3em; margin-bottom: 0.4em; color: #222; }
h3 { font-size: 13pt; margin-top: 1.1em; margin-bottom: 0.3em; color: #333; }
h4 { font-size: 12pt; margin-top: 1em; margin-bottom: 0.3em; font-style: italic; }
p { margin-bottom: 0.8em; text-align: justify; }
ul, ol { margin-bottom: 0.8em; padding-left: 2em; }
li { margin-bottom: 0.3em; }
blockquote {
    border-left: 3px solid #ccc;
    margin-left: 0;
    padding-left: 1em;
    color: #555;
    font-style: italic;
}
code {
    font-family: 'Courier New', monospace;
    font-size: 10pt;
    background: #f5f5f5;
    padding: 1px 4px;
}
pre {
    background: #f5f5f5;
    padding: 12px;
    font-size: 10pt;
    line-height: 1.4;
    white-space: pre-wrap;
    word-wrap: break-word;
}
pre code { background: none; padding: 0; }
table {
    width: 100%;
    border-collapse: collapse;
    margin: 1em 0;
    font-size: 11pt;
}
th, td {
    border: 1px solid #ccc;
    padding: 8px 10px;
    text-align: left;
}
th { background: #f0f0f0; font-weight: bold; }
hr { border: none; border-top: 1px solid #ddd; margin: 2em 0; }
"""

_REPORT_CSS = """
@page {
    size: A4;
    margin: 2cm;
}
body {
    font-family: Helvetica, Arial, sans-serif;
    font-size: 11pt;
    line-height: 1.5;
    color: #2d2d2d;
}
h1 { font-size: 24pt; margin-top: 1.2em; margin-bottom: 0.4em; color: #1a1a1a; font-weight: bold; }
h2 { font-size: 16pt; margin-top: 1.2em; margin-bottom: 0.4em; color: #333; border-bottom: 1px solid #eee; padding-bottom: 4px; }
h3 { font-size: 13pt; margin-top: 1em; margin-bottom: 0.3em; color: #444; }
h4 { font-size: 11pt; margin-top: 0.8em; margin-bottom: 0.3em; color: #555; font-weight: bold; }
p { margin-bottom: 0.7em; }
ul, ol { margin-bottom: 0.7em; padding-left: 1.8em; }
li { margin-bottom: 0.2em; }
blockquote {
    border-left: 4px solid #4f46e5;
    margin-left: 0;
    padding-left: 1em;
    color: #555;
}
code {
    font-family: 'Courier New', monospace;
    font-size: 9.5pt;
    background: #f8f8f8;
    padding: 2px 5px;
}
pre {
    background: #f8f8f8;
    padding: 14px;
    font-size: 9.5pt;
    line-height: 1.4;
    border: 1px solid #e8e8e8;
    white-space: pre-wrap;
    word-wrap: break-word;
}
pre code { background: none; padding: 0; border: none; }
table {
    width: 100%;
    border-collapse: collapse;
    margin: 1em 0;
    font-size: 10pt;
}
th, td {
    border: 1px solid #ddd;
    padding: 8px 12px;
    text-align: left;
}
th { background: #f5f5f5; font-weight: bold; }
hr { border: none; border-top: 2px solid #eee; margin: 1.5em 0; }
"""

_MINIMAL_CSS = """
@page {
    size: A4;
    margin: 1.5cm;
}
body {
    font-family: Helvetica, Arial, sans-serif;
    font-size: 10.5pt;
    line-height: 1.45;
    color: #333;
}
h1 { font-size: 18pt; margin-top: 1em; margin-bottom: 0.3em; }
h2 { font-size: 14pt; margin-top: 0.9em; margin-bottom: 0.3em; }
h3 { font-size: 12pt; margin-top: 0.8em; margin-bottom: 0.2em; }
p { margin-bottom: 0.5em; }
ul, ol { margin-bottom: 0.5em; padding-left: 1.5em; }
li { margin-bottom: 0.15em; }
code { font-family: monospace; font-size: 9.5pt; background: #f5f5f5; padding: 1px 3px; }
pre { background: #f5f5f5; padding: 10px; font-size: 9pt; white-space: pre-wrap; word-wrap: break-word; }
pre code { background: none; padding: 0; }
table { width: 100%; border-collapse: collapse; margin: 0.8em 0; font-size: 9.5pt; }
th, td { border: 1px solid #ddd; padding: 6px 8px; }
th { background: #f0f0f0; font-weight: bold; }
hr { border: none; border-top: 1px solid #ddd; margin: 1em 0; }
"""

_STYLE_CSS = {
    "academic": _ACADEMIC_CSS,
    "report": _REPORT_CSS,
    "minimal": _MINIMAL_CSS,
}


class DocumentGenerationService:
    """Generates PDF and DOCX documents from HTML content."""

    def __init__(self):
        pass

    async def generate_document(
        self,
        format: str,
        title: str,
        content: str,
        style: str = "academic",
        user_id: str | None = None,
    ) -> dict[str, Any]:
        """
        Generate a document in the specified format.

        Args:
            format: Document format ("pdf" or "docx")
            title: Document title
            content: HTML content to render (also handles markdown via conversion)
            style: Document style ("academic", "report", "minimal")
            user_id: User ID for path namespacing

        Returns:
            dict with keys: filename, url, size, format, content_type
        """
        if len(content) > MAX_CONTENT_LENGTH:
            content = content[:MAX_CONTENT_LENGTH]
            logger.warning(f"Content truncated to {MAX_CONTENT_LENGTH} chars for user {user_id}")

        format = format.lower().strip()
        if format not in ("pdf", "docx", "pptx"):
            raise ValueError(f"Unsupported format: {format}. Use 'pdf', 'docx', or 'pptx'.")

        # For pptx, content is a JSON string of slides; for pdf/docx it's HTML/markdown
        if format == "pptx":
            doc_bytes = self._generate_pptx(title, content, style)
            preview_html = self._build_pptx_preview_html(title, content, style)
        else:
            # Normalize content: if it's markdown, convert to HTML
            content = self._ensure_html(content)
            if format == "pdf":
                doc_bytes = self._generate_pdf(title, content, style)
            else:
                doc_bytes = self._generate_docx(title, content, style)
            preview_html = self._build_full_html(title, content, style)

        # Generate a unique filename
        safe_title = re.sub(r"[^\w\s-]", "", title)[:50].strip().replace(" ", "_")
        timestamp = datetime.now(UTC).strftime("%Y%m%d_%H%M%S")
        short_id = uuid.uuid4().hex[:6]
        filename = f"{safe_title}_{timestamp}_{short_id}.{format}"

        # Upload to storage
        from src.services.storage_service import storage_service

        storage_path = f"generated-docs/{user_id or 'anonymous'}"
        upload_result = await self._upload_bytes(storage_service, doc_bytes, filename, storage_path)

        # Also upload the styled HTML for in-app preview
        html_filename = f"{safe_title}_{timestamp}_{short_id}.html"
        preview_result = await self._upload_bytes(
            storage_service, preview_html.encode("utf-8"), html_filename, storage_path
        )

        return {
            "filename": filename,
            "url": upload_result["url"],
            "size": upload_result["size"],
            "format": format,
            "content_type": CONTENT_TYPES[format],
            "title": title,
            "preview_url": preview_result["url"],
        }

    def _ensure_html(self, content: str) -> str:
        """If content is markdown, convert to HTML. If already HTML, return as-is."""
        # Detect if content is already HTML (has block-level HTML tags)
        if re.search(r"<(h[1-6]|p|ul|ol|li|table|div|pre|blockquote)\b", content, re.IGNORECASE):
            return content

        # Content is markdown — convert to HTML
        return self._markdown_to_html(content)

    def _markdown_to_html(self, md: str) -> str:
        """Convert markdown to HTML for document rendering."""
        lines = md.split("\n")
        html_parts: list[str] = []
        in_code_block = False
        code_buffer: list[str] = []
        in_list = False
        list_type = ""
        paragraph_buffer: list[str] = []

        def flush_paragraph():
            if paragraph_buffer:
                text = " ".join(paragraph_buffer)
                text = self._inline_md_to_html(text)
                html_parts.append(f"<p>{text}</p>")
                paragraph_buffer.clear()

        def flush_list():
            nonlocal in_list, list_type
            if in_list:
                html_parts.append(f"</{list_type}>")
                in_list = False
                list_type = ""

        for line in lines:
            # Code blocks
            if line.strip().startswith("```"):
                if in_code_block:
                    html_parts.append(
                        f"<pre><code>{self._escape_html(chr(10).join(code_buffer))}</code></pre>"
                    )
                    code_buffer = []
                    in_code_block = False
                else:
                    flush_paragraph()
                    flush_list()
                    in_code_block = True
                continue

            if in_code_block:
                code_buffer.append(line)
                continue

            # Headings
            heading_match = re.match(r"^(#{1,6})\s+(.+)$", line)
            if heading_match:
                flush_paragraph()
                flush_list()
                level = len(heading_match.group(1))
                text = self._inline_md_to_html(heading_match.group(2))
                html_parts.append(f"<h{level}>{text}</h{level}>")
                continue

            # Horizontal rule
            if re.match(r"^(\-{3,}|\*{3,}|_{3,})\s*$", line):
                flush_paragraph()
                flush_list()
                html_parts.append("<hr>")
                continue

            # Bullet lists
            bullet_match = re.match(r"^(\s*)([-*+])\s+(.+)$", line)
            if bullet_match:
                flush_paragraph()
                if not in_list or list_type != "ul":
                    flush_list()
                    html_parts.append("<ul>")
                    in_list = True
                    list_type = "ul"
                text = self._inline_md_to_html(bullet_match.group(3))
                html_parts.append(f"<li>{text}</li>")
                continue

            # Numbered lists
            numbered_match = re.match(r"^(\s*)\d+\.\s+(.+)$", line)
            if numbered_match:
                flush_paragraph()
                if not in_list or list_type != "ol":
                    flush_list()
                    html_parts.append("<ol>")
                    in_list = True
                    list_type = "ol"
                text = self._inline_md_to_html(numbered_match.group(2))
                html_parts.append(f"<li>{text}</li>")
                continue

            # Blockquotes
            if line.startswith("> "):
                flush_paragraph()
                flush_list()
                text = self._inline_md_to_html(line[2:])
                html_parts.append(f"<blockquote><p>{text}</p></blockquote>")
                continue

            # Empty line
            if not line.strip():
                flush_paragraph()
                flush_list()
                continue

            # Regular text — accumulate for paragraph
            flush_list()
            paragraph_buffer.append(line.strip())

        # Flush remaining
        flush_paragraph()
        flush_list()
        if in_code_block and code_buffer:
            html_parts.append(
                f"<pre><code>{self._escape_html(chr(10).join(code_buffer))}</code></pre>"
            )

        return "\n".join(html_parts)

    def _inline_md_to_html(self, text: str) -> str:
        """Convert inline markdown (bold, italic, code, links) to HTML."""
        # Bold
        text = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", text)
        text = re.sub(r"__(.+?)__", r"<b>\1</b>", text)
        # Italic
        text = re.sub(r"\*(.+?)\*", r"<i>\1</i>", text)
        text = re.sub(r"_(.+?)_", r"<i>\1</i>", text)
        # Inline code
        text = re.sub(r"`(.+?)`", r"<code>\1</code>", text)
        # Links
        text = re.sub(r"\[(.+?)\]\((.+?)\)", r'<a href="\2">\1</a>', text)
        return text

    def _escape_html(self, text: str) -> str:
        """Escape HTML special characters."""
        return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    def _build_full_html(self, title: str, body_html: str, style: str) -> str:
        """Wrap body HTML with full document structure and CSS."""
        css = _STYLE_CSS.get(style, _ACADEMIC_CSS)
        return f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<title>{self._escape_html(title)}</title>
<style>{css}</style>
</head>
<body>
{body_html}
</body>
</html>"""

    def _generate_pdf(self, title: str, content: str, style: str) -> bytes:
        """Generate a PDF document from HTML content using xhtml2pdf."""
        from xhtml2pdf import pisa

        full_html = self._build_full_html(title, content, style)
        buffer = io.BytesIO()
        pisa_status = pisa.CreatePDF(io.StringIO(full_html), dest=buffer)

        if pisa_status.err:
            logger.error(f"xhtml2pdf error count: {pisa_status.err}")

        buffer.seek(0)
        return buffer.getvalue()

    def _generate_docx(self, title: str, content: str, style: str) -> bytes:
        """Generate a DOCX document from HTML content using htmldocx."""
        from docx import Document
        from docx.shared import Pt
        from htmldocx import HtmlToDocx

        doc = Document()
        parser = HtmlToDocx()

        # Add the HTML content
        parser.add_html_to_document(content, doc)

        # Save to bytes
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    def _generate_pptx(self, title: str, content: str, style: str) -> bytes:
        """Generate a richly styled PPTX from content with tables, shapes, and layouts."""
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Color palette
        PRIMARY = RGBColor(0x4F, 0x46, 0xE5)
        DARK_BG = RGBColor(0x1E, 0x1B, 0x4B)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
        TEXT_DARK = RGBColor(0x1E, 0x29, 0x3B)
        TEXT_GRAY = RGBColor(0x64, 0x74, 0x8B)
        ACCENT2 = RGBColor(0x7C, 0x3A, 0xED)
        ACCENT3 = RGBColor(0x06, 0xB6, 0xD4)

        slides_data = self._extract_slides_from_content(content, title)

        for i, slide_data in enumerate(slides_data):
            slide_title = slide_data.get("title", f"Slide {i + 1}")
            bullets = slide_data.get("bullets", [])
            subtitle = slide_data.get("subtitle", "")
            table_data = slide_data.get("table")

            layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(layout)

            if i == 0:
                # ═══ TITLE SLIDE ═══
                bg = slide.background.fill
                bg.solid()
                bg.fore_color.rgb = DARK_BG

                # Top accent bar
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.12)
                )
                bar.fill.solid()
                bar.fill.fore_color.rgb = PRIMARY
                bar.line.fill.background()

                # Left accent stripe
                acc = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(0.08), Inches(2.5)
                )
                acc.fill.solid()
                acc.fill.fore_color.rgb = PRIMARY
                acc.line.fill.background()

                # Title
                tb = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10), Inches(1.8))
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_title
                p.font.size = Pt(48)
                p.font.bold = True
                p.font.color.rgb = WHITE

                # Subtitle
                sub_text = subtitle or (bullets[0] if bullets else "")
                if sub_text:
                    tb2 = slide.shapes.add_textbox(Inches(1.2), Inches(4.2), Inches(10), Inches(1))
                    tf2 = tb2.text_frame
                    tf2.word_wrap = True
                    p2 = tf2.paragraphs[0]
                    p2.text = sub_text
                    p2.font.size = Pt(20)
                    p2.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)

                # Decorative circles
                for cx, clr in [(Inches(10), ACCENT2), (Inches(11.2), ACCENT3)]:
                    c = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL, cx, Inches(5.5), Inches(1.5), Inches(1.5)
                    )
                    c.fill.solid()
                    c.fill.fore_color.rgb = clr
                    c.line.fill.background()

            elif table_data and isinstance(table_data, dict):
                # ═══ TABLE SLIDE ═══
                self._pptx_add_header(slide, slide_title, prs.slide_width, PRIMARY, WHITE)
                headers = table_data.get("headers", [])
                rows = table_data.get("rows", [])
                if headers and rows:
                    n_rows = len(rows) + 1
                    n_cols = len(headers)
                    tbl_w = min(prs.slide_width - Inches(1.6), Inches(n_cols * 2.5))
                    tbl_left = (prs.slide_width - tbl_w) // 2
                    shape = slide.shapes.add_table(
                        n_rows, n_cols, tbl_left, Inches(2.0), tbl_w, Inches(4.5)
                    )
                    tbl = shape.table
                    tbl.first_row = True
                    for ci, h in enumerate(headers):
                        cell = tbl.cell(0, ci)
                        cell.text = str(h)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = PRIMARY
                        for para in cell.text_frame.paragraphs:
                            para.font.color.rgb = WHITE
                            para.font.size = Pt(14)
                            para.font.bold = True
                    for ri, row in enumerate(rows):
                        for ci, val in enumerate(row[:n_cols]):
                            cell = tbl.cell(ri + 1, ci)
                            cell.text = str(val)
                            for para in cell.text_frame.paragraphs:
                                para.font.size = Pt(13)
                                para.font.color.rgb = TEXT_DARK
                            if ri % 2 == 1:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                # ═══ CONTENT SLIDE ═══
                self._pptx_add_header(slide, slide_title, prs.slide_width, PRIMARY, WHITE)
                if len(bullets) > 6:
                    mid = len(bullets) // 2
                    self._pptx_add_bullets(
                        slide,
                        bullets[:mid],
                        Inches(0.8),
                        Inches(2.0),
                        Inches(5.5),
                        TEXT_DARK,
                        PRIMARY,
                    )
                    self._pptx_add_bullets(
                        slide,
                        bullets[mid:],
                        Inches(6.8),
                        Inches(2.0),
                        Inches(5.5),
                        TEXT_DARK,
                        ACCENT2,
                    )
                elif bullets:
                    self._pptx_add_bullets(
                        slide, bullets, Inches(0.8), Inches(2.0), Inches(11.5), TEXT_DARK, PRIMARY
                    )

            # Slide number
            if i > 0:
                nb = slide.shapes.add_textbox(Inches(12.4), Inches(7.0), Inches(0.7), Inches(0.4))
                np = nb.text_frame.paragraphs[0]
                np.text = str(i)
                np.font.size = Pt(10)
                np.font.color.rgb = TEXT_GRAY
                np.alignment = PP_ALIGN.RIGHT

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    def _pptx_add_header(self, slide, title_text, slide_width, color, text_color):
        """Add styled header bar to a slide."""
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt

        hdr = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), slide_width, Inches(1.5)
        )
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11.5), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.LEFT

    def _pptx_add_bullets(self, slide, bullets, left, top, width, text_color, dot_color):
        """Add styled bullet list with dot shapes."""
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt

        y = top
        for bullet in bullets:
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, left, y + Inches(0.12), Inches(0.15), Inches(0.15)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = dot_color
            dot.line.fill.background()
            tb = slide.shapes.add_textbox(left + Inches(0.35), y, width - Inches(0.35), Inches(0.6))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(bullet)
            p.font.size = Pt(18)
            p.font.color.rgb = text_color
            y += Inches(0.7)

    def _extract_slides_from_content(self, content: str, title: str) -> list[dict]:
        """Extract slide data from content (HTML sections, JSON, or plain text)."""
        import json

        # Try JSON first
        try:
            data = json.loads(content)
            if isinstance(data, list) and len(data) > 0:
                return data
            if isinstance(data, dict) and "slides" in data:
                return data["slides"]
        except (json.JSONDecodeError, TypeError):
            pass

        # Try HTML <section> format
        if "<section" in content:
            return self._parse_sections_to_slides(content)

        # Fallback: use the legacy parser
        return self._parse_pptx_content(content, title)

    def _parse_sections_to_slides(self, html: str) -> list[dict]:
        """Parse HTML <section> tags into slide data for PPTX, including tables."""
        slides: list[dict] = []
        # Split by section tags
        sections = re.split(r"<section[^>]*>", html)

        for section in sections:
            if not section.strip():
                continue
            # Remove closing tag
            section = re.sub(r"</section>.*", "", section, flags=re.DOTALL)

            # Extract title from h1 or h2
            title_match = re.search(r"<h[12][^>]*>(.*?)</h[12]>", section, re.DOTALL)
            slide_title = ""
            if title_match:
                slide_title = re.sub(r"<[^>]+>", "", title_match.group(1)).strip()

            # Extract subtitle from first <p> after title
            subtitle = ""
            subtitle_match = re.search(r"</h[12]>\s*<p[^>]*>(.*?)</p>", section, re.DOTALL)
            if subtitle_match:
                subtitle = re.sub(r"<[^>]+>", "", subtitle_match.group(1)).strip()

            # Extract table data if present
            table_data = None
            table_match = re.search(
                r"<table[^>]*>(.*?)</table>", section, re.DOTALL | re.IGNORECASE
            )
            if table_match:
                table_html = table_match.group(1)
                headers: list[str] = []
                rows: list[list[str]] = []
                # Extract headers from <th> tags
                for th in re.finditer(r"<th[^>]*>(.*?)</th>", table_html, re.DOTALL):
                    headers.append(re.sub(r"<[^>]+>", "", th.group(1)).strip())
                # Extract rows from <td> tags grouped by <tr>
                for tr in re.finditer(r"<tr[^>]*>(.*?)</tr>", table_html, re.DOTALL):
                    tr_content = tr.group(1)
                    if "<th" in tr_content:
                        continue  # Skip header row
                    row: list[str] = []
                    for td in re.finditer(r"<td[^>]*>(.*?)</td>", tr_content, re.DOTALL):
                        row.append(re.sub(r"<[^>]+>", "", td.group(1)).strip())
                    if row:
                        rows.append(row)
                if headers and rows:
                    table_data = {"headers": headers, "rows": rows}

            # Extract bullet points from <li> tags
            bullets: list[str] = []
            for li_match in re.finditer(r"<li[^>]*>(.*?)</li>", section, re.DOTALL):
                text = re.sub(r"<[^>]+>", "", li_match.group(1)).strip()
                if text:
                    bullets.append(text)

            # If no bullets and no table, extract <p> content as bullets
            if not bullets and not table_data:
                p_tags = re.findall(r"<p[^>]*>(.*?)</p>", section, re.DOTALL)
                for p in p_tags:
                    text = re.sub(r"<[^>]+>", "", p).strip()
                    if text and text != subtitle:
                        bullets.append(text)

            if slide_title or bullets or table_data:
                slide: dict = {"title": slide_title or "Untitled", "bullets": bullets}
                if subtitle:
                    slide["subtitle"] = subtitle
                if table_data:
                    slide["table"] = table_data
                slides.append(slide)

        if not slides:
            slides = [{"title": "Presentation", "subtitle": "", "bullets": []}]

        return slides

    def _parse_pptx_content(self, content: str, title: str) -> list[dict]:
        """Parse content into slide data. Supports JSON array or plain text fallback."""
        import json

        # Try JSON first
        try:
            data = json.loads(content)
            if isinstance(data, list) and len(data) > 0:
                return data
            if isinstance(data, dict) and "slides" in data:
                return data["slides"]
        except (json.JSONDecodeError, TypeError):
            pass

        # Fallback: parse HTML/text content into slides by splitting on headings
        slides: list[dict] = []
        current_slide: dict | None = None

        # Strip HTML tags for text extraction
        text = re.sub(r"<[^>]+>", "\n", content)
        text = re.sub(r"\n{3,}", "\n\n", text).strip()
        lines = text.split("\n")

        for line in lines:
            line = line.strip()
            if not line:
                continue

            # Detect heading-like lines (short, no period at end, capitalized)
            is_heading = (
                len(line) < 100
                and not line.endswith(".")
                and (line[0].isupper() or line[0].isdigit())
                and not line.startswith("-")
                and not line.startswith("*")
            )

            if is_heading and (not current_slide or len(current_slide.get("bullets", [])) > 0):
                if current_slide:
                    slides.append(current_slide)
                current_slide = {"title": line, "bullets": []}
            elif current_slide:
                # Clean bullet markers
                clean = re.sub(r"^[-*•]\s*", "", line)
                if clean:
                    current_slide["bullets"].append(clean)
            else:
                current_slide = {"title": title, "bullets": [line]}

        if current_slide:
            slides.append(current_slide)

        # If no slides parsed, create a single title slide
        if not slides:
            slides = [{"title": title, "subtitle": "Generated presentation", "bullets": []}]

        return slides

    def _build_pptx_preview_html(self, title: str, content: str, style: str) -> str:
        """Build an HTML preview for a PPTX presentation from HTML section slides."""
        # Check if content uses <section> tags (new rich HTML format)
        if "<section" in content:
            body = content
        else:
            # Legacy JSON format fallback
            import json

            slides_data = self._parse_pptx_content(content, title)
            slides_parts = []
            for i, slide in enumerate(slides_data):
                slide_title = slide.get("title", f"Slide {i + 1}")
                bullets = slide.get("bullets", [])
                subtitle = slide.get("subtitle", "")
                bullets_html = ""
                if bullets:
                    items = "".join(f"<li>{self._escape_html(b)}</li>" for b in bullets)
                    bullets_html = f"<ul>{items}</ul>"
                subtitle_html = f"<p>{self._escape_html(subtitle)}</p>" if subtitle else ""
                slides_parts.append(
                    f"<section><h2>{self._escape_html(slide_title)}</h2>"
                    f"{subtitle_html}{bullets_html}</section>"
                )
            body = "\n".join(slides_parts)

        return self._wrap_presentation_html(title, body)

    def _wrap_presentation_html(self, title: str, body: str) -> str:
        """Wrap slide sections with reveal.js for interactive slideshow + standalone CSS fallback."""
        css = """
* { box-sizing: border-box; }
.reveal section {
    padding: 40px;
    text-align: left;
}
.reveal h1 {
    font-size: 2.2em;
    font-weight: 800;
    margin-bottom: 0.3em;
    color: #111;
}
.reveal h2 {
    font-size: 1.6em;
    font-weight: 700;
    margin-bottom: 0.5em;
    color: #1a1a1a;
    border-bottom: 3px solid #4f46e5;
    padding-bottom: 8px;
    display: inline-block;
}
.reveal h3 {
    font-size: 1.2em;
    font-weight: 600;
    margin: 0.8em 0 0.4em;
    color: #333;
}
.reveal p {
    font-size: 0.95em;
    margin-bottom: 0.6em;
    color: #444;
    line-height: 1.6;
}
.reveal ul, .reveal ol {
    display: block;
    padding-left: 1.5em;
    margin: 0.5em 0;
}
.reveal li {
    font-size: 0.9em;
    margin-bottom: 0.5em;
    color: #333;
    line-height: 1.5;
}
.reveal table {
    width: 100%;
    border-collapse: collapse;
    margin: 1em 0;
    font-size: 0.85em;
}
.reveal th {
    background: #4f46e5;
    color: white;
    padding: 10px 14px;
    text-align: left;
    font-weight: 600;
}
.reveal td {
    padding: 10px 14px;
    border-bottom: 1px solid #eee;
    color: #333;
}
.reveal tr:nth-child(even) td { background: #f8f8ff; }
.reveal blockquote {
    border-left: 4px solid #4f46e5;
    padding: 12px 16px;
    margin: 1em 0;
    background: #f5f3ff;
    border-radius: 0 8px 8px 0;
    font-style: italic;
    color: #4a4a6a;
    width: 90%;
}
.reveal blockquote cite {
    display: block;
    margin-top: 8px;
    font-style: normal;
    font-size: 0.85em;
    color: #666;
    font-weight: 600;
}
.reveal svg {
    display: block;
    margin: 1em auto;
    max-width: 100%;
}
.columns {
    display: flex;
    gap: 20px;
    margin: 1em 0;
}
.columns > div {
    flex: 1;
    background: #f8fafc;
    padding: 16px;
    border-radius: 10px;
    border: 1px solid #e2e8f0;
}
.columns > div h3 { margin-top: 0; }
.stat {
    display: inline-flex;
    flex-direction: column;
    align-items: center;
    background: #f0f0ff;
    padding: 20px 28px;
    border-radius: 12px;
    margin: 8px 12px 8px 0;
}
.stat .number {
    font-size: 2.5em;
    font-weight: 800;
    color: #4f46e5;
    line-height: 1;
}
.stat .label {
    font-size: 0.75em;
    color: #666;
    margin-top: 6px;
    font-weight: 600;
    text-transform: uppercase;
    letter-spacing: 0.5px;
}
.highlight {
    background: linear-gradient(135deg, #4f46e5 0%, #7c3aed 100%);
    color: white;
    padding: 20px 24px;
    border-radius: 12px;
    margin: 1em 0;
}
.highlight h3, .highlight p, .highlight li { color: white !important; }
.highlight ul { padding-left: 20px; }
.timeline { margin: 1em 0; }
.timeline .event {
    padding: 12px 0 12px 24px;
    border-left: 3px solid #4f46e5;
    margin-bottom: 6px;
    position: relative;
}
.timeline .event::before {
    content: '';
    position: absolute;
    left: -7px;
    top: 16px;
    width: 11px;
    height: 11px;
    border-radius: 50%;
    background: #4f46e5;
}
.timeline .event b { color: #4f46e5; }
.timeline .event p { margin: 4px 0 0; font-size: 0.9em; color: #555; }
"""
        return f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{self._escape_html(title)}</title>
<link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@5.1.0/dist/reveal.css">
<link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@5.1.0/dist/theme/white.css">
<style>{css}</style>
</head>
<body>
<div class="reveal">
<div class="slides">
{body}
</div>
</div>
<script src="https://cdn.jsdelivr.net/npm/reveal.js@5.1.0/dist/reveal.js"></script>
<script>
Reveal.initialize({{
    hash: true,
    slideNumber: true,
    controls: true,
    progress: true,
    center: false,
    transition: 'slide',
    width: 1280,
    height: 720,
}});
</script>
</body>
</html>"""

    async def _upload_bytes(
        self, storage_service: Any, content: bytes, filename: str, path: str
    ) -> dict[str, Any]:
        """Upload raw bytes to BunnyCDN storage."""
        import httpx

        if not storage_service.api_key or not storage_service.storage_zone:
            raise RuntimeError("Storage configuration is missing.")

        upload_path = f"{path.strip('/')}/{filename}"
        upload_url = f"{storage_service.base_url}/{upload_path}"

        headers = {
            "AccessKey": storage_service.api_key,
            "Content-Type": "application/octet-stream",
        }

        async with httpx.AsyncClient(timeout=60.0) as client:
            response = await client.put(upload_url, headers=headers, content=content)
            if response.status_code != 201:
                raise RuntimeError(f"Upload failed: {response.status_code} - {response.text}")

        # Build public URL
        if storage_service.public_url_base:
            public_url = f"{storage_service.public_url_base}/{upload_path}"
        else:
            public_url = f"https://{storage_service.cdn_hostname}/{upload_path}"

        return {"filename": filename, "url": public_url, "size": len(content)}


# Module-level singleton
document_generation_service = DocumentGenerationService()
